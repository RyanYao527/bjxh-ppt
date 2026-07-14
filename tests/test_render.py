"""Tests for render.py utility functions and error paths.

Full integration tests (render_page with real BJXH layouts) require the
master template .pptx which is not in the repository.  These tests cover
the functions that *can* be exercised with a default python-pptx template.
"""

from __future__ import annotations

import sys
from pathlib import Path

import pytest

sys.path.insert(0, str(Path(__file__).parent.parent / "scripts"))
from render import (  # noqa: E402
    PLACEHOLDER_MAP,
    find_layout,
    find_placeholder,
    render_page,
    set_placeholder_text,
)
from parse import PageSpec  # noqa: E402
from pptx import Presentation  # noqa: E402


@pytest.fixture(scope="module")
def prs() -> Presentation:
    """A fresh default Presentation with built-in layouts."""
    return Presentation()


@pytest.fixture(scope="module")
def slide(prs: Presentation) -> object:
    """A slide from the default 'Title Slide' layout, if available."""
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if "title" in layout.name.lower():
                return prs.slides.add_slide(layout)
    # Fallback: first available layout
    layout = prs.slide_masters[0].slide_layouts[0]
    return prs.slides.add_slide(layout)


# -- find_layout -----------------------------------------------------------


def test_find_layout_returns_none_for_bogus_name(prs: Presentation) -> None:
    assert find_layout(prs, "this layout does not exist") is None


def test_find_layout_finds_default_layout(prs: Presentation) -> None:
    # The default pptx template always has a blank layout or similar
    layout = prs.slide_masters[0].slide_layouts[0]
    found = find_layout(prs, layout.name)
    assert found is not None
    assert found.name == layout.name


# -- find_placeholder ------------------------------------------------------


def test_find_placeholder_returns_none_for_bogus_idx(slide: object) -> None:
    assert find_placeholder(slide, 999) is None


def test_find_placeholder_finds_valid_idx(slide: object) -> None:
    # Get a real placeholder idx from the slide
    if slide.placeholders:
        real_idx = slide.placeholders[0].placeholder_format.idx
        ph = find_placeholder(slide, real_idx)
        assert ph is not None
        assert ph.placeholder_format.idx == real_idx


# -- set_placeholder_text --------------------------------------------------


def test_set_placeholder_text_writes_and_clears(slide: object) -> None:
    if not slide.placeholders:
        pytest.skip("Slide has no placeholders")
    ph = slide.placeholders[0]
    if not ph.has_text_frame:
        pytest.skip("Placeholder has no text frame")

    result = set_placeholder_text(ph, "Test Title", size_pt=18, bold=True)
    assert result is True
    assert "Test Title" in ph.text_frame.text

    # Overwrite
    set_placeholder_text(ph, "New Title", size_pt=20, bold=False)
    text = ph.text_frame.text
    assert "New Title" in text
    assert "Test Title" not in text


def test_set_placeholder_text_returns_false_for_none() -> None:
    assert set_placeholder_text(None, "text") is False


# -- render_page error paths -----------------------------------------------


def test_render_page_unknown_layout_raises(prs: Presentation) -> None:
    spec = PageSpec(kind="content", title="Test", layout="不存在的版式")
    with pytest.raises(SystemExit, match="layout.*not found"):
        render_page(prs, spec, template_path="/nonexistent.pptx")


def test_render_page_layout_not_in_map_raises(prs: Presentation) -> None:
    # Find a layout that exists in the default template but is NOT in
    # PLACEHOLDER_MAP (e.g. the default 'Blank' or 'Title Slide' layout).
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name not in PLACEHOLDER_MAP:
                spec = PageSpec(
                    kind="content", title="Test", layout=layout.name
                )
                with pytest.raises(SystemExit, match="not in PLACEHOLDER_MAP"):
                    render_page(prs, spec, template_path="/nonexistent.pptx")
                return
    pytest.skip("All default layouts are in PLACEHOLDER_MAP")
