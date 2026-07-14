"""Tests for layout_utils.py — font sizing, layout suggestion, text helpers."""

from __future__ import annotations

import sys
from pathlib import Path

import pytest

sys.path.insert(0, str(Path(__file__).parent.parent / "scripts"))
from layout_utils import (  # noqa: E402
    calc_safe_font_size,
    clear_unused_placeholders,
    count_cjk,
    estimated_char_width_pt,
    extract_segment_title,
    suggest_layout,
    truncate_title,
)
from pptx import Presentation  # noqa: E402


# -- count_cjk -------------------------------------------------------------


@pytest.mark.parametrize("text,expected", [
    ("纯中文文本", 5),
    ("Hello World", 0),
    ("混合ABC中文", 4),
    ("", 0),
    ("emoji😀中文", 3),  # emoji codepoint > 0x2E80 counts as CJK
])
def test_count_cjk(text: str, expected: int) -> None:
    assert count_cjk(text) == expected


# -- estimated_char_width_pt -----------------------------------------------


def test_char_width_pure_cjk() -> None:
    assert estimated_char_width_pt("中文") == pytest.approx(1.0, abs=0.01)


def test_char_width_pure_ascii() -> None:
    assert estimated_char_width_pt("hello") == pytest.approx(0.55, abs=0.01)


def test_char_width_mixed() -> None:
    w = estimated_char_width_pt("中a")  # 1.0*1 + 0.55*1 / 2 = 0.775
    assert 0.7 < w < 0.85


# -- calc_safe_font_size ---------------------------------------------------


class MockPh:
    """Minimal placeholder mock with width/height in EMU."""
    def __init__(self, w_in: float, h_in: float):
        self.width = int(w_in * 914400)
        self.height = int(h_in * 914400)


def test_safe_font_cjk_short_fits_max() -> None:
    ph = MockPh(3.87, 0.54)
    # Short CJK title fits at max_pt
    assert calc_safe_font_size(ph, "核心判断", max_pt=20) == 20


def test_safe_font_long_text_scales_down() -> None:
    ph = MockPh(3.87, 0.54)
    # Long text should reduce below max_pt
    result = calc_safe_font_size(ph, "审计行业AI应用分析报告与核心判断", max_pt=20)
    assert result < 20


def test_safe_font_empty_text() -> None:
    ph = MockPh(3.87, 0.54)
    assert calc_safe_font_size(ph, "", max_pt=20) == 20


def test_safe_font_never_below_min() -> None:
    ph = MockPh(1.0, 0.2)  # Tiny box
    assert calc_safe_font_size(ph, "非常长的中文标题文本内容", max_pt=20, min_pt=8) == 8


def test_safe_font_ascii_long() -> None:
    ph = MockPh(5.0, 1.0)
    # ASCII chars are narrower, should fit more text at larger size
    result = calc_safe_font_size(ph, "A" * 60, max_pt=20)
    assert result >= 14  # ASCII fits better than CJK


# -- suggest_layout --------------------------------------------------------


def test_suggest_layout_short_bullets() -> None:
    layout = suggest_layout(["要点一", "要点二"])
    assert layout in ("无图分段-3项", "无图分段-4项", "无图分段-5项",
                      "标题页-空白", "2_标题页-空白",
                      "有图分段式-16", "有图分段式-8")


def test_suggest_layout_four_items_prefers_4() -> None:
    layout = suggest_layout(["a", "b", "c", "d"])
    assert "4项" in layout or layout in ("有图分段式-16", "有图分段式-8")


def test_suggest_layout_five_items_prefers_5() -> None:
    layout = suggest_layout(["a", "b", "c", "d", "e"])
    assert "5项" in layout or layout in ("有图分段式-16", "有图分段式-8")


def test_suggest_layout_data_triggers_chart() -> None:
    # Contains data keywords like % and 增长
    layout = suggest_layout(["收入增长20%", "成本下降5%", "利润率提升"])
    assert layout in ("图表-1", "图表-2",
                      "无图分段-3项", "有图分段式-16", "有图分段式-8")


def test_suggest_layout_excludes_recent() -> None:
    used = frozenset(["无图分段-3项", "无图分段-4项", "无图分段-5项"])
    layout = suggest_layout(["a", "b"], used_layouts=used)
    # Should pick text or image-text since all structured are used
    assert layout in ("标题页-空白", "2_标题页-空白",
                      "有图分段式-16", "有图分段式-8")


def test_suggest_layout_empty_bullets() -> None:
    layout = suggest_layout([])
    assert layout in ("标题页-空白", "2_标题页-空白",
                      "无图分段-3项", "无图分段-4项", "无图分段-5项")


# -- truncate_title --------------------------------------------------------


def test_truncate_short_title_unchanged() -> None:
    display, note = truncate_title("核心判断", max_chars=14)
    assert display == "核心判断"
    assert note == ""


def test_truncate_at_colon() -> None:
    display, note = truncate_title("审计行业AI应用：从效率工具到战略变量", max_chars=14)
    assert display == "审计行业AI应用"
    assert note != ""


def test_truncate_no_separator() -> None:
    display, note = truncate_title("这是一个非常长的标题没有任何分隔符", max_chars=14)
    assert len(display) == 14
    assert note != ""


def test_truncate_exact_max_chars() -> None:
    display, note = truncate_title("一二三四五六七八九十一二三四", max_chars=14)
    assert len(display) == 14


# -- extract_segment_title -------------------------------------------------


def test_extract_chinese_colon() -> None:
    assert extract_segment_title("成本结构变革：AI将证据获取的边际成本压至趋零") == "成本结构变革"


def test_extract_english_colon() -> None:
    assert extract_segment_title("Key: value pair of text") == "Key"


def test_extract_short_text() -> None:
    assert extract_segment_title("收费压力") == "收费压力"


def test_extract_parentheses_removal() -> None:
    result = extract_segment_title("第一阶段(0-6月)基础建设方案")
    assert "0-6月" not in result
    assert len(result) <= 8


def test_extract_long_fallback() -> None:
    result = extract_segment_title("这是非常长的文本没有任何分隔符可以用")
    assert len(result) == 8


@pytest.fixture(scope="module")
def slide() -> object:
    """A slide from a default Presentation."""
    prs = Presentation()
    layout = prs.slide_masters[0].slide_layouts[0]
    return prs.slides.add_slide(layout)


# -- clear_unused_placeholders ---------------------------------------------


def test_clear_unused_placeholder(slide: object) -> None:
    """Write a space to an unused placeholder, verify it's cleared."""
    if not slide.placeholders:
        pytest.skip("No placeholders")
    # Pick a placeholder, don't add it to used_idxs, call clear
    ph = slide.placeholders[0]
    if not ph.has_text_frame:
        pytest.skip("Placeholder has no text frame")
    # Set some template text
    ph.text_frame.text = "单击此处添加文本"
    clear_unused_placeholders(slide, used_idx_set=set())
    # Should now be cleared to a space
    assert ph.text_frame.text.strip() == ""


def test_clear_used_placeholder_preserved(slide: object) -> None:
    """Used placeholder should NOT be cleared."""
    if not slide.placeholders:
        pytest.skip("No placeholders")
    ph = slide.placeholders[0]
    if not ph.has_text_frame:
        pytest.skip("Placeholder has no text frame")
    ph.text_frame.text = "Important Content"
    used = {ph.placeholder_format.idx}
    clear_unused_placeholders(slide, used_idx_set=used)
    assert "Important" in ph.text_frame.text
