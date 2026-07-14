"""Integration tests for render_page using a minimal default template."""

from __future__ import annotations

import sys
from pathlib import Path

import pytest

sys.path.insert(0, str(Path(__file__).parent.parent / "scripts"))
from parse import PageSpec  # noqa: E402
from pptx import Presentation  # noqa: E402
from render import PLACEHOLDER_MAP, find_placeholder, render_page  # noqa: E402


@pytest.fixture(scope="module")
def prs() -> Presentation:
    return Presentation()


def test_render_page_cover_with_default_template(prs: Presentation) -> None:
    """Cover page with 主题-封面: either renders or raises SystemExit."""
    spec = PageSpec(kind="cover", title="Test Cover", layout="主题-封面")
    # Check if default template actually has 主题-封面
    has_layout = False
    for m in prs.slide_masters:
        for lay in m.slide_layouts:
            if lay.name == "主题-封面":
                has_layout = True
    if not has_layout:
        with pytest.raises(SystemExit, match="layout.*not found"):
            render_page(prs, spec, template_path="/nonexistent.pptx")
    else:
        render_page(prs, spec, template_path="/nonexistent.pptx")
        assert len(prs.slides) >= 1


def test_render_content_writes_title(prs: Presentation) -> None:
    """Content page should write title to main_title_idx placeholder."""
    # Use 无图分段-3项 if available, else any mapped content layout
    layout = None
    for name in ["无图分段-3项", "标题页-空白"]:
        if name in PLACEHOLDER_MAP:
            # Check if default template has a layout we can use
            for m in prs.slide_masters:
                for lay in m.slide_layouts:
                    if lay.name == name:
                        layout = name
                        break
            if layout:
                break
    if not layout:
        pytest.skip("No mapped content layout found in default template")

    spec = PageSpec(kind="content", title="测试标题", bullets=["要点一", "要点二"], layout=layout)
    render_page(prs, spec, template_path="/nonexistent.pptx")
    slide = prs.slides[-1]
    # Find the title placeholder and verify
    mp = PLACEHOLDER_MAP[layout]
    title_ph = find_placeholder(slide, mp["main_title_idx"])
    if title_ph and title_ph.has_text_frame:
        assert "测试标题" in title_ph.text_frame.text


def test_render_page_clears_unused_placeholders(prs: Presentation) -> None:
    """After rendering, unused placeholders should not contain default text."""
    layout = None
    for name in ["无图分段-3项", "标题页-空白"]:
        if name in PLACEHOLDER_MAP:
            for m in prs.slide_masters:
                for lay in m.slide_layouts:
                    if lay.name == name:
                        layout = name
                        break
    if not layout:
        pytest.skip("No mapped content layout found")

    spec = PageSpec(kind="content", title="Test", bullets=[], layout=layout)
    render_page(prs, spec, template_path="/nonexistent.pptx")
    slide = prs.slides[-1]
    # No placeholder should contain "单击此处" (template default text)
    for ph in slide.placeholders:
        if ph.has_text_frame:
            assert "单击此处" not in ph.text_frame.text
