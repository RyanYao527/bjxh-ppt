"""Add custom diagrams to the AI协同 PPT content slides."""
from pptx import Presentation
from pptx.enum.dml import MSO_THEME_COLOR_INDEX
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

RED = MSO_THEME_COLOR_INDEX.ACCENT_3
GRAY = MSO_THEME_COLOR_INDEX.ACCENT_2
DARK = MSO_THEME_COLOR_INDEX.TEXT_1
WHITE = MSO_THEME_COLOR_INDEX.BACKGROUND_1


def set_text(shape, text: str, *, size_pt: int = 14, bold: bool = False, color=MSO_THEME_COLOR_INDEX.TEXT_1):
    tf = shape.text_frame
    tf.text = ""
    p = tf.paragraphs[0]
    p.line_spacing = 1.2
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.theme_color = color
    p.alignment = 1  # center


def add_card(slide, x, y, w, h, title, body, *, fill_color=RED, title_size=16, body_size=12):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.theme_color = fill_color
    card.line.fill.background()
    set_text(card, f"{title}\n{body}", size_pt=body_size, bold=False, color=WHITE)
    return card


def add_arrow(slide, x, y, w, h, color=RED):
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    arrow.fill.solid()
    arrow.fill.fore_color.theme_color = color
    arrow.line.fill.background()
    return arrow


def add_label(slide, x, y, w, h, text, *, size_pt=14, bold=False, color=MSO_THEME_COLOR_INDEX.TEXT_1):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text(tb, text, size_pt=size_pt, bold=bold, color=color)
    return tb


def enhance_slide_3(slide):
    """传统模式 vs AI协同 —— 左右对比卡片 + 中央箭头."""
    # Left card: traditional mode
    add_card(slide, 0.8, 3.0, 4.5, 3.0,
             "传统模式",
             "串行推进\n任务频繁切换\n受个人时间精力限制",
             fill_color=GRAY, title_size=18, body_size=14)

    # Central arrow
    add_arrow(slide, 5.5, 4.0, 2.2, 0.9, color=RED)

    # Right card: AI协同模式
    add_card(slide, 8.0, 3.0, 4.5, 3.0,
             "AI 协同模式",
             "多任务并行\nAgent 各司其职\n我聚焦决策与整合",
             fill_color=RED, title_size=18, body_size=14)

    # Bottom summary
    add_label(slide, 2.0, 6.3, 9.5, 0.5,
              "核心变化：我负责决策，AI 负责执行",
              size_pt=16, bold=True, color=RED)


def enhance_slide_4(slide):
    """Agent 分工体系 —— 2×2 能力矩阵."""
    # Title row already at top; add matrix below
    matrix = [
        ("开发执行", "Codex GUI / CLI\n工程实现 · 长文档", 0.9, 2.5, RED),
        ("研究分析", "Gemini / DeepSeek / ChatGPT\n深度研究 · 推理验证", 6.9, 2.5, GRAY),
        ("中文处理", "Kimi / GLM\n资料搜集 · 表达优化", 0.9, 4.3, GRAY),
        ("演示输出", "MiniMax\nPPT 制作 · 视觉化", 6.9, 4.3, RED),
    ]
    for title, body, x, y, color in matrix:
        add_card(slide, x, y, 5.4, 1.6, title, body, fill_color=color, title_size=16, body_size=13)

    # Center cross lines (subtle)
    hline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.4), Inches(2.5), Inches(0.04), Inches(3.4))
    hline.fill.solid(); hline.fill.fore_color.theme_color = MSO_THEME_COLOR_INDEX.ACCENT_1; hline.line.fill.background()
    vline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(4.2), Inches(12.5), Inches(0.04))
    vline.fill.solid(); vline.fill.fore_color.theme_color = MSO_THEME_COLOR_INDEX.ACCENT_1; vline.line.fill.background()

    # Bottom principle
    add_label(slide, 2.0, 6.3, 9.5, 0.4,
              "没有最强 Agent，只有最适合当前任务的 Agent",
              size_pt=15, bold=True, color=DARK)


def enhance_slide_5(slide):
    """效率之变 —— 公式 + 前后对比."""
    # Formula row
    y = 3.0
    h = 0.9
    # 1人
    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), Inches(y), Inches(1.4), Inches(h))
    oval.fill.solid(); oval.fill.fore_color.theme_color = GRAY; oval.line.fill.background()
    set_text(oval, "1人", size_pt=18, bold=True, color=WHITE)

    add_label(slide, 2.5, y + 0.25, 0.6, 0.5, "+", size_pt=24, bold=True)

    # 4~6 Agents
    rect1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.2), Inches(y), Inches(2.4), Inches(h))
    rect1.fill.solid(); rect1.fill.fore_color.theme_color = RED; rect1.line.fill.background()
    set_text(rect1, "4~6 个 Agent", size_pt=16, bold=True, color=WHITE)

    add_label(slide, 5.7, y + 0.25, 0.8, 0.5, "=", size_pt=24, bold=True)

    # 4~5倍产出
    rect2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.6), Inches(y), Inches(3.2), Inches(h))
    rect2.fill.solid(); rect2.fill.fore_color.theme_color = RED; rect2.line.fill.background()
    set_text(rect2, "4~5 倍产出", size_pt=22, bold=True, color=WHITE)

    # Before/After
    add_card(slide, 1.0, 4.3, 4.8, 2.2,
             "过去",
             "我完成工作\n时间是核心瓶颈\n产出依赖个人能力",
             fill_color=GRAY, title_size=16, body_size=13)

    add_arrow(slide, 5.9, 5.0, 1.6, 0.8, color=RED)

    add_card(slide, 7.7, 4.3, 4.8, 2.2,
             "现在",
             "我组织 AI 完成工作\n多任务同时推进\n我专注方向与质量",
             fill_color=RED, title_size=16, body_size=13)


def main():
    path = r"C:\Users\rori9\AI协同工作实践分享_v2.pptx"
    prs = Presentation(path)
    slides = list(prs.slides)
    # slides[2]=page3, [3]=page4, [4]=page5 (0-based)
    enhance_slide_3(slides[2])
    enhance_slide_4(slides[3])
    enhance_slide_5(slides[4])
    prs.save(path)
    print(f"Enhanced: {path}")


if __name__ == "__main__":
    main()
