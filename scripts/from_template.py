"""
from_template.py — Create a clean .pptx by selecting slide layouts from the
Beijing Xinghua master template.

Usage:
    python from_template.py <template.pptx> <out.pptx> <layout1,layout2,...>

Examples:
    # Standard opening: cover + TOC + section + content + closing
    python from_template.py master.pptx out.pptx "主题-封面,主题-目录页,标题页-空白,文字模板1,主题-封底页"

The output file inherits the template's theme (colors, fonts, master, layouts)
and is ready to be opened in PowerPoint for content entry.
"""
import argparse
import sys
from pathlib import Path

from pptx import Presentation

from shared import apply_minimal_formatting  # noqa: F401 — re-exported for compatibility


# 2026 兴华模板 — 标准 12 个常用版式（详见 SKILL.md §6.1）
STANDARD_OPENING = "主题-封面,主题-目录页,标题页-空白"
STANDARD_CLOSING = "主题-封底页"


def parse_layout_list(raw: str) -> list[str]:
    return [s.strip() for s in raw.split(",") if s.strip()]


def find_layout(prs: Presentation, name: str):
    """Find a slide layout by exact name across all slide masters."""
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == name:
                return layout
    return None


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("template", help="path to 北京兴华模板.pptx")
    parser.add_argument("output", help="output .pptx path")
    parser.add_argument(
        "layouts",
        help='comma-separated layout names, e.g. "主题-封面,主题-目录页,文字模板1"',
    )
    args = parser.parse_args()

    template_path = Path(args.template)
    if not template_path.exists():
        print(f"ERROR: template not found: {template_path}", file=sys.stderr)
        return 1

    prs = Presentation(str(template_path))

    layout_names = parse_layout_list(args.layouts)
    missing: list[str] = []
    new_slides = []

    for name in layout_names:
        layout = find_layout(prs, name)
        if layout is None:
            missing.append(name)
            continue
        slide = prs.slides.add_slide(layout)
        apply_minimal_formatting(slide)
        new_slides.append((name, slide))

    if missing:
        print(
            f"ERROR: {len(missing)} layout(s) not found in template: {missing}",
            file=sys.stderr,
        )
        print(
            "Run dump_template_spec.py to see all available layouts.",
            file=sys.stderr,
        )
        return 2

    out_path = Path(args.output)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))

    print(f"OK: wrote {out_path} with {len(new_slides)} slide(s)")
    for i, (name, _slide) in enumerate(new_slides, 1):
        print(f"  slide {i}: {name}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
