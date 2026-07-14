"""
layout_utils.py — Layout helpers for bjxh-ppt render pipeline.

Provides dynamic font sizing (prevents text overflow) and placeholder
cleanup (removes "Click here to add text" residue).
"""

from __future__ import annotations

import sys as _sys
from typing import Any

from pptx.util import Pt

from shared import emu_to_inches


def count_cjk(text: str) -> int:
    """Return the number of CJK / wide characters in *text*."""
    return sum(1 for c in text if ord(c) > 0x2E80)


def estimated_char_width_pt(text: str) -> float:
    """Rough per-character width in points, accounting for CJK vs ASCII.

    CJK characters are roughly 1.0× font-size wide; ASCII ~0.55×.
    Returns a weighted average for the given text.
    """
    if not text:
        return 0.55
    cjk = count_cjk(text)
    ascii_count = len(text) - cjk
    return (1.0 * cjk + 0.55 * ascii_count) / len(text)


def calc_safe_font_size(
    ph: Any,
    text: str,
    *,
    max_pt: int = 20,
    min_pt: int = 8,
    line_spacing: float = 1.3,
) -> int:
    """Find the largest font size (pt) that fits *text* in *ph*.

    Steps down from *max_pt* until the text fits both width and height,
    bottoming out at *min_pt*.
    """
    if not text:
        return max_pt

    box_w = emu_to_inches(ph.width) if ph.width else 13.3
    box_h = emu_to_inches(ph.height) if ph.height else 7.5
    char_w_ratio = estimated_char_width_pt(text)

    for size in range(max_pt, min_pt - 1, -1):
        char_width_in = (size / 72) * char_w_ratio
        chars_per_line = max(1, int(box_w / char_width_in))
        lines_needed = -(-len(text) // chars_per_line)  # ceil division
        line_height_in = (size / 72) * line_spacing
        lines_fit = max(1, int(box_h / line_height_in))
        if lines_needed <= lines_fit:
            return size
    return min_pt


# Content-aware layout categories (all names must match PLACEHOLDER_MAP).
# Only layouts in PLACEHOLDER_MAP are eligible for auto-selection.
_STRUCTURED = ["无图分段-3项", "无图分段-4项", "无图分段-5项"]
_TEXT = ["标题页-空白", "2_标题页-空白"]
_CHART = ["图表-1", "图表-2"]

# Keywords that hint at content type
_DATA_KEYWORDS = ["%", "增长", "下降", "占比", "同比", "环比", "数据", "比例",
                  "亿", "万", "倍", "个百分点", "统计", "趋势"]
_PROCESS_KEYWORDS = ["步骤", "流程", "阶段", "首先", "然后", "最后", "环节",
                     "先后", "第一步", "第二步"]


def suggest_layout(bullets: list[str], *, used_layouts: frozenset[str] | None = None) -> str:
    """Suggest an appropriate layout for a content page based on its bullets.

    Heuristics (in priority order):
    1. Data-heavy bullets → chart layout
    2. Process/step bullets → logic layout
    3. Short bullets (≤3 items) → structured card layout
    4. Everything else → rotated through structured variants

    *used_layouts* is a set of layouts already used nearby; the function
    avoids repeating the same layout.
    """
    used = used_layouts or frozenset()
    text = " ".join(bullets)

    # Data-heavy bullets → try chart layout first
    if any(kw in text for kw in _DATA_KEYWORDS) and len(bullets) <= 6:
        for c in _CHART:
            if c not in used:
                return c

    # Assemble all available non-structured variants for rotation
    all_candidates = list(_STRUCTURED)
    # Insert text layouts occasionally for variety
    if len(bullets) <= 3:
        all_candidates = _TEXT + all_candidates

    # Exclude recently-used layouts
    candidates = [lay for lay in all_candidates if lay not in used]
    if not candidates:
        candidates = all_candidates

    # Prefer match on bullet count for structured layouts
    n = len(bullets)
    structured_preferred = [lay for lay in candidates if lay in _STRUCTURED]
    text_preferred = [lay for lay in candidates if lay in _TEXT]

    if n <= 3:
        # Short pages: alternate between text and structured
        if text_preferred and recent_layouts_structured(used) >= 2:
            return text_preferred[0]
        ordered = [lay for lay in candidates if "3项" in lay or lay in _TEXT]
    elif n == 4:
        ordered = [lay for lay in candidates if "4项" in lay]
    elif n >= 5:
        ordered = [lay for lay in candidates if "5项" in lay]
    else:
        ordered = []

    fallback = ordered or structured_preferred or text_preferred or candidates
    return fallback[0]


def extract_segment_title(bullet: str, max_chars: int = 8) -> str:
    """Extract a short title from a bullet string for segment card headers.

    Priority: separator-split prefix → first N CJK chars → raw truncation.
    Avoids producing garbled fragments like '第一阶段(0' or '成本结构变'.
    """
    # 1) Try splitting on colon / Chinese colon — use the prefix
    for sep in ["：", ":"]:
        if sep in bullet:
            prefix = bullet.split(sep, 1)[0].strip()
            if 2 <= len(prefix) <= max_chars + 2:
                return prefix
            # If prefix is too long, take first N chars of it
            if len(prefix) > max_chars:
                return prefix[:max_chars]

    # 2) Try splitting on comma / Chinese comma / period — prefix
    for sep in ["，", ",", "。", "."]:
        if sep in bullet:
            prefix = bullet.split(sep, 1)[0].strip()
            if 2 <= len(prefix) <= max_chars:
                return prefix

    # 3) Try parentheses → drop parenthetical content for cleaner title
    import re
    cleaned = re.sub(r"[（(][^)）]*[)）]", "", bullet)
    if 2 <= len(cleaned) <= max_chars:
        return cleaned.strip()
    if len(cleaned) > max_chars:
        return cleaned[:max_chars].strip()

    # 4) Fallback: first N chars
    return bullet[:max_chars]


def recent_layouts_structured(used: frozenset[str]) -> int:
    """Count how many recently-used layouts are structured variants."""
    return sum(1 for lay in used if lay in _STRUCTURED)


def clear_unused_placeholders(slide: Any, used_idx_set: set[int]) -> None:
    """Clear text from any placeholder whose idx is NOT in *used_idx_set*.

    This prevents "单击此处添加文本" (Click to add text) template residue
    from appearing in the final .pptx when PowerPoint renders unfilled
    placeholders.
    """
    for ph in slide.placeholders:
        idx: int = ph.placeholder_format.idx
        if idx not in used_idx_set and ph.has_text_frame:
            # Write a single space so PowerPoint doesn't show placeholder prompt
            tf = ph.text_frame
            tf.text = " "
