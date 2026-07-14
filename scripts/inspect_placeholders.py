"""
inspect_placeholders.py — Dump placeholder metadata for every layout in a
template .pptx.  Used to verify / regenerate PLACEHOLDER_MAP entries in
render.py when the master template changes.

Usage:
    python inspect_placeholders.py <template.pptx> [--layout "主题-封面"]

Output (JSON to stdout):
    For each layout (or a single --layout): name, index, and a list of
    placeholders with idx, type, name, position (inches), and size (inches).
"""

from __future__ import annotations

import json
import sys
from pathlib import Path

from pptx import Presentation

from shared import emu_to_inches


def dump_layouts(template_path: str, filter_name: str | None = None) -> list[dict]:
    prs = Presentation(template_path)
    results: list[dict] = []

    for m_idx, master in enumerate(prs.slide_masters):
        for l_idx, layout in enumerate(master.slide_layouts):
            name: str = layout.name
            if filter_name and name != filter_name:
                continue

            placeholders: list[dict] = []
            for ph in layout.placeholders:
                phf = ph.placeholder_format
                placeholders.append(
                    {
                        "idx": phf.idx,
                        "type": str(phf.type),
                        "name": ph.name,
                        "shape_type": str(ph.shape_type),
                        "left_in": emu_to_inches(ph.left) if ph.left is not None else None,
                        "top_in": emu_to_inches(ph.top) if ph.top is not None else None,
                        "width_in": emu_to_inches(ph.width) if ph.width is not None else None,
                        "height_in": emu_to_inches(ph.height) if ph.height is not None else None,
                        "has_text_frame": ph.has_text_frame,
                    }
                )

            results.append(
                {
                    "master_index": m_idx,
                    "layout_index": l_idx,
                    "layout_name": name,
                    "placeholder_count": len(placeholders),
                    "placeholders": placeholders,
                }
            )

    return results


def main() -> int:
    if len(sys.argv) < 2:
        print(__doc__)
        return 1

    template_path: str = sys.argv[1]
    if not Path(template_path).exists():
        print(f"ERROR: template not found: {template_path}", file=sys.stderr)
        return 1

    filter_name: str | None = None
    if len(sys.argv) >= 4 and sys.argv[2] == "--layout":
        filter_name = sys.argv[3]

    results: list[dict] = dump_layouts(template_path, filter_name)
    print(
        json.dumps(
            {"source": template_path, "layout_count": len(results), "layouts": results},
            ensure_ascii=False,
            indent=2,
        )
    )
    return 0


if __name__ == "__main__":
    sys.exit(main())
