"""
qa.py — Beijing Xinghua PPT QA: enforce the 9 must-check items in SKILL.md §8.

Usage:
    python qa.py <file.pptx>

Exit code 0 = PASS, non-zero = FAIL with detailed issues printed.
"""
from __future__ import annotations

import re
import sys
from collections import Counter
from collections.abc import Generator
from pathlib import Path
from typing import Any

from pptx import Presentation


CANVAS_W: int = 12192000
CANVAS_H: int = 6858000

# §2 禁止字体
FORBIDDEN_FONTS: set[str] = {
    "宋体",
    "SimSun",
    "黑体",
    "SimHei",
    "楷体",
    "KaiTi",
    "隶书",
    "LiSu",
    "华文宋体",
    "华文黑体",
    "华文楷体",
    "华文隶书",
    "STSong",
    "STHeiti",
    "STKaiti",
    "Times New Roman",
    "Calibri",
    "Cambria",
    "Georgia",
}

# §4.2 模板编辑提示色
EDIT_NOTE_GREY: tuple[int, int, int] = (0x44, 0x54, 0x69)  # #445469

# §7 占位符孤立字母 — matches isolated B/J/H/X characters (old template residue)
# as well as consecutive runs or dot/space-separated pairs.
BJH_PATTERN: re.Pattern[str] = re.compile(
    r"(?:^|\s|[，。：；、])[BJHX](?:$|\s|[，。：；、])"  # single isolated char
    r"|"
    r"[BJHX]{2,}"  # two or more consecutive (e.g. "BJH")
    r"|"
    r"[BJHX][.\s][BJHX]"  # dot- or space-separated pair (e.g. "B.J", "B H")
)


def iter_text_runs(slide: Any) -> Generator[tuple[Any, Any, Any], None, None]:
    """Yield (shape, paragraph, run) for every non-empty run on a slide."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text and run.text.strip():
                    yield shape, para, run


def check_canvas(prs: Presentation, issues: list[str]) -> None:
    if prs.slide_width != CANVAS_W or prs.slide_height != CANVAS_H:
        issues.append(
            f"[CANVAS] wrong size: {prs.slide_width}x{prs.slide_height} EMU "
            f"(expected {CANVAS_W}x{CANVAS_H}, i.e. 13.333\"x7.5\")"
        )


def check_fonts_and_size(prs: Presentation, issues: list[str]) -> None:
    for s_idx, slide in enumerate(prs.slides, 1):
        for _shape, _para, run in iter_text_runs(slide):
            text: str = run.text.strip()
            # §2 禁止字体
            if run.font.name and run.font.name in FORBIDDEN_FONTS:
                issues.append(
                    f"[FONT] slide {s_idx}: forbidden font '{run.font.name}' "
                    f"in '{text[:40]}'"
                )
            # §9.2 font.name / font.size / font.color must be explicit
            if run.font.name is None:
                issues.append(
                    f"[FONT-NONE] slide {s_idx}: font.name is None in '{text[:40]}'"
                )
            if run.font.size is None:
                issues.append(
                    f"[SIZE-NONE] slide {s_idx}: font.size is None in '{text[:40]}'"
                )


def check_layout_variety(prs: Presentation, issues: list[str]) -> None:
    layouts: list[str] = [s.slide_layout.name for s in prs.slides]
    counter: Counter[str] = Counter(layouts)
    distinct: int = len(counter)
    total: int = len(layouts)
    if total >= 3 and distinct < 3:
        issues.append(
            f"[LAYOUT] only {distinct} distinct layout(s) across {total} slides; "
            f"spec requires at least 3. distribution: {dict(counter)}"
        )


def check_essential_placeholders_filled(prs: Presentation, issues: list[str]) -> None:
    """Per layout, the main title placeholder must be non-empty. For
    `无图分段-N项` layouts, we count non-empty segment placeholders: at
    least 1 segment must be filled (the layout was chosen, so SOMETHING
    should be in it). 4 项 with 4 empty segments is a real bug.

    This is the check that catches the catastrophic case where python-pptx
    wrote the title into the wrong placeholder and left the actual title
    placeholder empty.
    """
    MIN_FILL: dict[str, tuple[int, int]] = {
        "主题-封面":   (10, 0),
        "标题页-空白": (11, 0),
        "无图分段-3项": (11, 1),
        "无图分段-4项": (11, 1),
        "无图分段-5项": (11, 1),
        "主题-封底页":  (10, 0),
    }
    SEG_BODY_IDXS: dict[str, list[int]] = {
        "无图分段-3项": [18, 19, 20],
        "无图分段-4项": [28, 30, 32, 34],
        "无图分段-5项": [22, 24, 26, 28, 30],
    }

    for s_idx, slide in enumerate(prs.slides, 1):
        layout_name: str = slide.slide_layout.name
        rule: tuple[int, int] | None = MIN_FILL.get(layout_name)
        if rule is None:
            continue
        title_idx, min_filled = rule

        idx_text: dict[int, str] = {}
        for ph in slide.placeholders:
            if ph.has_text_frame:
                idx_text[ph.placeholder_format.idx] = ph.text_frame.text.strip()

        title_text: str = idx_text.get(title_idx, "")
        if not title_text:
            issues.append(
                f"[PLACEHOLDER-EMPTY] slide {s_idx} (layout '{layout_name}'): "
                f"main title idx={title_idx} is empty"
            )

        if layout_name in SEG_BODY_IDXS:
            filled: int = sum(1 for i in SEG_BODY_IDXS[layout_name] if idx_text.get(i))
            if filled < min_filled:
                issues.append(
                    f"[NO-CONTENT] slide {s_idx} (layout '{layout_name}'): "
                    f"{filled} segment bodies filled (expected at least {min_filled})"
                )


def check_bjh_residue(prs: Presentation, issues: list[str]) -> None:
    for s_idx, slide in enumerate(prs.slides, 1):
        for _shape, _para, run in iter_text_runs(slide):
            text: str = run.text
            if BJH_PATTERN.search(text):
                issues.append(
                    f"[BJH] slide {s_idx}: B/J/X/H placeholder residue in '{text[:60]}'"
                )


def check_edit_grey_color(prs: Presentation, issues: list[str]) -> None:
    for s_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not run.text or not run.text.strip():
                        continue
                    try:
                        color = run.font.color
                        if color is None:
                            continue
                        # Only RGB-type colors can have the forbidden #445469.
                        # Theme colors (SCHEME type) and inherited colors are fine.
                        if not hasattr(color, "rgb") or color.rgb is None:
                            continue
                        rgb = color.rgb
                        if (
                            rgb[0] == EDIT_NOTE_GREY[0]
                            and rgb[1] == EDIT_NOTE_GREY[1]
                            and rgb[2] == EDIT_NOTE_GREY[2]
                        ):
                            issues.append(
                                f"[GREY-NOTE] slide {s_idx}: edit-note grey "
                                f"#445469 found in '{run.text[:40]}'"
                            )
                    except (AttributeError, TypeError) as exc:
                        print(
                            f"[WARN] slide {s_idx}: could not inspect color of "
                            f"'{run.text[:40]}': {exc}",
                            file=sys.stderr,
                        )


def check_visual_element_per_slide(prs: Presentation, issues: list[str]) -> None:
    """A page is "too thin" if it has very few shapes AND they're all text.

    Cover / closing / chapter-title / TOC pages are exempt — they are
    intentionally minimal. Content pages should have >= 4 shapes (text + the
    template's built-in decorative blocks, which are also text frames in
    python-pptx's eyes).
    """
    EXEMPT_LAYOUT_KEYWORDS: tuple[str, ...] = ("封面", "封底", "标题页", "目录", "自定义版式", "过渡页")
    MIN_CONTENT_SHAPES: int = 4
    for s_idx, slide in enumerate(prs.slides, 1):
        layout_name: str = slide.slide_layout.name
        if any(k in layout_name for k in EXEMPT_LAYOUT_KEYWORDS):
            continue
        n_shapes: int = len(slide.shapes)
        if n_shapes < MIN_CONTENT_SHAPES:
            issues.append(
                f"[THIN-PAGE] slide {s_idx} (layout '{layout_name}'): "
                f"only {n_shapes} shape(s); content pages need >= {MIN_CONTENT_SHAPES}. "
                f"Add a chart/table/picture or pick a richer layout."
            )


def main() -> int:
    if len(sys.argv) != 2:
        print(__doc__)
        return 1

    path = Path(sys.argv[1])
    if not path.exists():
        print(f"ERROR: file not found: {path}", file=sys.stderr)
        return 1

    prs = Presentation(str(path))
    issues: list[str] = []

    check_canvas(prs, issues)
    check_fonts_and_size(prs, issues)
    check_layout_variety(prs, issues)
    check_essential_placeholders_filled(prs, issues)
    check_bjh_residue(prs, issues)
    check_edit_grey_color(prs, issues)
    check_visual_element_per_slide(prs, issues)

    if issues:
        print(f"FAIL: {len(issues)} issue(s) in {path.name}")
        for i, issue in enumerate(issues, 1):
            print(f"  {i}. {issue}")
        return 2

    print(f"PASS: {path.name} ({len(prs.slides)} slides) meets 北京兴华 PPT 规范")
    return 0


if __name__ == "__main__":
    sys.exit(main())
