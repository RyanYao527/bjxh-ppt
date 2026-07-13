"""
shared.py — utilities shared across bjxh-ppt scripts.

Formerly, apply_minimal_formatting was duplicated in from_template.py
and imported via a brittle sys.path hack in from_outline.py.
"""

from pptx.util import Pt


def apply_minimal_formatting(slide) -> None:
    """Apply minimal but spec-compliant formatting to every text run.

    - font.name = 'Microsoft YaHei'
    - font.size = Pt(16) (body default) unless already set

    This is a safety net; most placeholders inherit from layouts.
    """
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font.name is None:
                    run.font.name = "Microsoft YaHei"
                if run.font.size is None:
                    run.font.size = Pt(16)
