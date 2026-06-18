"""
from_outline.py — Convert a Markdown outline to a Beijing Xinghua .pptx.

Outline syntax (deliberately minimal — see examples/audit_demo.md for full sample):

    # 主标题 (H1)        — 整份演示的标题，仅第一行使用 → 主题-封面
    ## 章节名 (H2)        — 章节封面 → 标题页-空白
    ### 页面标题 (H3)     — 一张内容页 → 默认版式（可指定）
    - 要点                 — 页内要点
    - 要点 2
    > layout: 文字模板1   — 显式指定当前 H3 章节的版式（可选，覆盖默认）
    > note: 备注文字       — 添加演讲者备注（可选）

Usage:
    python from_outline.py outline.md out.pptx [template.pptx]
        template.pptx 默认为 SKILL.md §12 列出的标准路径

版式自动选择规则（见 SKILL.md §6.1）:
    H1 (1 个)              → 主题-封面
    第一个 H2 之后插入     → 主题-目录页（可选：out_no_toc 关闭）
    每个 H2                → 标题页-空白（章节封面）
    每个 H3                → 文字模板1（默认）/ > layout: 覆盖
    末尾                    → 主题-封底页（可选：out_no_closing 关闭）

如果内容不足以触发默认版式，会用 `>` 指令显式指定。
"""
import argparse
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

try:
    from from_template import apply_minimal_formatting
except ImportError:
    # allow running from any cwd
    sys.path.insert(0, str(Path(__file__).parent))
    from from_template import apply_minimal_formatting  # type: ignore


DEFAULT_TEMPLATE = (
    r"C:\工作\04-总结与报告\2026年工作\2026合伙人大会\北京兴华模板.pptx"
)

# H3 缺省使用的版式（最常见、最易读）
DEFAULT_CONTENT_LAYOUT = "文字模板1"


# ---------------- markdown parsing ----------------

H1_RE = re.compile(r"^#\s+(.+?)\s*$")
H2_RE = re.compile(r"^##\s+(.+?)\s*$")
H3_RE = re.compile(r"^###\s+(.+?)\s*$")
UL_RE = re.compile(r"^[-*]\s+(.+?)\s*$")
DIRECTIVE_RE = re.compile(r"^>\s*(\w+):\s*(.+?)\s*$")


class PageSpec:
    """One page = one H3 (or an H2 chapter cover, or the cover/closing)."""

    def __init__(self, kind: str, title: str, layout: str, bullets: list[str], note: str = ""):
        self.kind = kind  # 'cover' | 'toc' | 'chapter' | 'content' | 'closing'
        self.title = title
        self.layout = layout
        self.bullets = bullets
        self.note = note

    def __repr__(self) -> str:
        return f"PageSpec({self.kind}, {self.title!r}, layout={self.layout!r}, n={len(self.bullets)})"


def parse_outline(text: str, *, add_toc: bool = True, add_closing: bool = True) -> list[PageSpec]:
    """Parse a markdown outline into a list of PageSpec objects."""
    pages: list[PageSpec] = []
    current_chapter: str | None = None
    current: PageSpec | None = None
    cover_added = False

    def flush() -> None:
        nonlocal current
        if current is not None:
            pages.append(current)
            current = None

    for raw_line in text.splitlines():
        line = raw_line.rstrip()
        if not line.strip():
            continue

        m = H1_RE.match(line)
        if m and not cover_added:
            flush()
            current = PageSpec(
                kind="cover",
                title=m.group(1).strip(),
                layout="主题-封面",
                bullets=[],
            )
            cover_added = True
            continue

        m = H2_RE.match(line)
        if m:
            flush()
            current_chapter = m.group(1).strip()
            current = PageSpec(
                kind="chapter",
                title=current_chapter,
                layout="标题页-空白",
                bullets=[],
            )
            continue

        m = H3_RE.match(line)
        if m:
            flush()
            current = PageSpec(
                kind="content",
                title=m.group(1).strip(),
                layout=DEFAULT_CONTENT_LAYOUT,
                bullets=[],
            )
            continue

        m = UL_RE.match(line)
        if m and current is not None:
            current.bullets.append(m.group(1).strip())
            continue

        m = DIRECTIVE_RE.match(line)
        if m and current is not None:
            key, val = m.group(1).lower(), m.group(2).strip()
            if key == "layout":
                current.layout = val
            elif key == "note":
                current.note = val
            continue

    flush()

    # 装饰：插入目录页 + 封底页
    if add_toc and any(p.kind == "chapter" for p in pages):
        # 收集所有 H2 章节标题作为目录项
        chapter_titles = [p.title for p in pages if p.kind == "chapter"]
        toc = PageSpec(
            kind="toc",
            title="目录",
            layout="主题-目录页",
            bullets=chapter_titles,
        )
        # 插到封面之后、第一章之前
        insert_at = 1 if pages and pages[0].kind == "cover" else 0
        pages.insert(insert_at, toc)

    if add_closing and pages and pages[-1].kind != "closing":
        pages.append(
            PageSpec(
                kind="closing",
                title="北京兴华集团",   # 封底使用公司名(模板 自定义版式 layout111 设计)
                layout="自定义版式",
                bullets=[],
            )
        )

    return pages


# ---------------- rendering ----------------

def find_layout(prs: Presentation, name: str):
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == name:
                return layout
    return None


def find_placeholder(slide, ph_idx: int):
    """Find a placeholder by exact idx. No guessing by type — see PLACEHOLDER_MAP."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == ph_idx:
            return ph
    return None


def set_placeholder_text(ph, text: str, *, size_pt: int = 16, bold: bool = False) -> bool:
    """Write `text` into placeholder, preserving the existing first run's formatting
    where possible. Returns True on success."""
    if ph is None or not ph.has_text_frame:
        return False
    tf = ph.text_frame
    tf.text = ""  # clear existing
    p = tf.paragraphs[0]
    p.line_spacing = 1.3
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    if run.font.color and run.font.color.type is None:
        from pptx.dml.color import RGBColor
        run.font.color.rgb = RGBColor(0x1F, 0x1F, 0x1F)
    return True


# 事实依据来自 inspect_placeholders 输出。每个 layout 都指明"哪个 idx 是什么角色"。
# 维护原则:模板变了就重跑 inspect_placeholders.py,不要靠猜。
PLACEHOLDER_MAP: dict[str, dict] = {
    "主题-封面": {
        "main_title_idx": 10,   # 文本占位符 8
        "subtitle_idx":  11,   # 文本占位符 10(英文副标题)
    },
    "主题-目录页": {
        # 事实(来自 inspect_placeholders 输出,2026-06-18 核验):
        #   4 个"章节标题"占位符 idx = [12, 21, 23, 25]   (左列,从上到下)
        #   4 个"英文小标题"占位符 idx = [20, 22, 24, 26] (紧贴每个章节标题上方的小字)
        #   4 个"对应页码"占位符 idx = [27, 29, 31, 33]   (右列,从上到下)
        #   4 个"小页码"占位符 idx = [28, 30, 32, 34]     (与页码同列的占位)
        # 只填章节标题(12,21,23,25)+ 页码(27,29,31,33);小标题/小页码留空即可。
        # 最多 4 章节,超出 4 个的需要换更复杂的目录版式。
        "item_idxs":  [12, 21, 23, 25],   # 章节中文标题(只填前 4 个)
        "item_sub_idxs": [20, 22, 24, 26],  # 章节英文小标题(可空)
        "page_idxs":  [27, 29, 31, 33],   # 对应页码
    },
    "标题页-空白": {
        "main_title_idx": 11,   # 文本占位符 2(一级标题)
        "subtitle_idx":  12,   # 文本占位符 6
    },
    "无图分段-3项": {
        "main_title_idx": 11,   # 文本占位符 2(页面标题)
        "subtitle_idx":  12,   # 文本占位符 6
        "seg_title_idxs": [15, 16, 17],  # 3 个分段小标题(idx 15-17)
        "seg_body_idxs":  [18, 19, 20],  # 3 个分段正文(idx 18-20)
    },
    "无图分段-4项": {
        "main_title_idx": 11,
        "subtitle_idx":  12,
        "seg_title_idxs": [27, 29, 31, 33],
        "seg_body_idxs":  [28, 30, 32, 34],
    },
    "无图分段-5项": {
        "main_title_idx": 11,
        "subtitle_idx":  12,
        "seg_title_idxs": [21, 23, 25, 27, 29],
        "seg_body_idxs":  [22, 24, 26, 28, 30],
    },
    "主题-封底页": {
        # 实测(2026-06-18):
        #   idx=10  文本占位符 "THANKS"  pos=(5.14, 2.99)in  size=3.04×0.89in
        #   idx=4   页码占位符 (SLIDE_NUMBER) 不写
        # 模板自带装饰(模板层,不需代码触碰):
        #   - 全幅背景图(13.33×7.50in)
        #   - 4 个装饰形状(红/白波浪)
        #   - "诚信铸就品牌 专业创造价值" 14pt 文本框
        "main_title_idx":   10,   # 唯一要写的占位符(主标题)
        "slide_number_idx": 4,    # 页码占位符(由 PowerPoint 自动填,代码不写)
    },
    "自定义版式": {
        # 实测(2026-06-18):
        #   idx=11  文本占位符 "北京兴华集团"  pos=(3.94, 3.01)in  size=5.44×0.89in
        #         默认 40pt Microsoft YaHei bold
        #   idx=10  页码占位符 (SLIDE_NUMBER)  模板自动填
        # 模板自带装饰(模板层,不需代码触碰):
        #   - 北京建筑背景(图片 2,13.33×7.50in)
        #   - 顶部右上角红色波浪装饰
        #   - 左上角"北京兴华 XH GLOBAL"logo
        #   - "诚信铸就品牌 专业创造价值" 14pt 文本框 (3.44, 3.68)in
        #   - 左下角联系方式区(电话/传真/地址)
        # **重要：不要动模板的背景/装饰/logo/联系方式,只填 idx=11 主标题即可**
        "main_title_idx":   11,   # 唯一要写的占位符(主标题,默认 40pt bold)
        "slide_number_idx": 10,   # 页码占位符
    },
}


def render_page(prs: Presentation, spec: PageSpec) -> None:
    layout = find_layout(prs, spec.layout)
    if layout is None:
        raise SystemExit(
            f"ERROR: layout '{spec.layout}' not found. "
            f"Check spelling against template_spec.json."
        )
    if spec.layout not in PLACEHOLDER_MAP:
        raise SystemExit(
            f"ERROR: layout '{spec.layout}' not in PLACEHOLDER_MAP. "
            f"Supported: {list(PLACEHOLDER_MAP)}"
        )
    slide = prs.slides.add_slide(layout)
    apply_minimal_formatting(slide)
    mp = PLACEHOLDER_MAP[spec.layout]

    if spec.kind == "cover":
        # 封面（主题-封面 = layout39）实测规范（2026-06-18 v3 修订）：
        #
        # layout 自带装饰（不需代码补）：
        #   - "BEIJING XINGHUA" 灰色水印  pos=(2.78, 2.20)in  size=8.03×1.01in
        #         54pt Arial Black 浅色填充 + softEdge 模糊效果
        #   - 左上角 logo（图片） pos=(0.94, 0.82)in  size=1.98×0.51in
        #   - 顶部右上角红色波浪 + 底部红色波浪
        #
        # 主标题占位符（idx=10）：
        #   - 模板默认 pos=(2.40, 2.62)in  size=8.54×1.18in
        #   - **必须 66pt**（用户规范 2026-06-18 v3 修订）
        #   - 66pt 文字在 1.18in 框内（66/72=0.92in 文字高），
        #     与水印框 y=2.20-3.21in 在 y=2.74-3.21in 区段自然视觉重叠 0.47in
        #   - 注：旧版 28pt 文字只占 0.39in 高，居中后位于 y=3.02-3.41in，
        #     完全脱离水印范围（顶部 2.46in），所以"调大到 66pt 才会重叠"
        #
        # 副标题占位符（idx=11）：
        #   - 模板默认 pos=(0.51, 6.10)in  size=4.04×1.39in
        #   - **必须 18pt**（用户规范 2026-06-18 v3 修订）
        #   - 模板默认内容"BEIJING XINGHUA CERTIFIED PUBLIC ACCOUNTANTS" 20pt 折行 3 行
        #   - 用户要"BEIJING XINGHUA  GROUP"（不写也行，留空用模板默认）
        #   - 注：旧版 28pt 文字宽 8.66in > 框宽 4.04in，强制换行成 2 行；改 18pt 后单行
        size = 66
        bold = True
        set_placeholder_text(
            find_placeholder(slide, mp["main_title_idx"]), spec.title, size_pt=size, bold=bold
        )
        # 副标题（英文机构名）：idx=11 写"BEIJING XINGHUA GROUP" 18pt 单行
        # 模板默认 3 行"BEIJING XINGHUA CERTIFIED PUBLIC ACCOUNTANTS" 20pt → 必须替换
        if mp.get("subtitle_idx") is not None:
            sub_ph = find_placeholder(slide, mp["subtitle_idx"])
            if sub_ph is not None and sub_ph.has_text_frame:
                # 重写副标题文本为"BEIJING XINGHUA GROUP" 18pt
                tf = sub_ph.text_frame
                tf.clear()  # 清空模板默认 3 行
                p = tf.paragraphs[0]
                p.alignment = None  # 左对齐
                run = p.add_run()
                run.text = "BEIJING XINGHUA  GROUP"
                run.font.name = "Arial"
                run.font.size = Pt(18)
                # 注：18pt "BEIJING XINGHUA  GROUP" ~3.5in 宽 < 框宽 4.04in → 单行
                # 旧版 28pt 文字 ~5.4in 宽 > 4.04in → 强制换行成 2 行（错误）

    elif spec.kind == "toc":
        # 目录:把 spec.bullets 当作"章节标题",依次填到 4 个 idx
        # 模板实际只有 4 个章节位置(12,21,23,25),超过 4 章节的需要换版式
        for i, idx in enumerate(mp["item_idxs"]):
            text = spec.bullets[i] if i < len(spec.bullets) else ""
            set_placeholder_text(find_placeholder(slide, idx), text, size_pt=18, bold=False)
        # 4 个页码占位符:自动算"第几页"(只填 N,不写"01"格式)
        # 简单起见:填 1, 2, 3, ...(用户在 PPT 里改成实际页码)
        for i, idx in enumerate(mp["page_idxs"]):
            if i < len(spec.bullets):
                set_placeholder_text(find_placeholder(slide, idx), str(i + 1), size_pt=18, bold=False)

    elif spec.kind == "chapter":
        # 章节标题页 — 模板框 3.87×0.54 寸,32pt 装不下 12+ 字,降到 24pt
        set_placeholder_text(
            find_placeholder(slide, mp["main_title_idx"]), spec.title, size_pt=24, bold=True
        )
        set_placeholder_text(
            find_placeholder(slide, mp["subtitle_idx"]), spec.note or "", size_pt=14, bold=False
        )

    elif spec.kind == "content":
        # 内容页:主页标题 + 副标题 + (如果有结构化分段版式)3 个分段
        # 模板框 3.87×0.54 寸,24pt 是单行极限,长标题会被裁
        set_placeholder_text(
            find_placeholder(slide, mp["main_title_idx"]), spec.title, size_pt=20, bold=True
        )
        if mp.get("subtitle_idx") is not None:
            # 副标题(可选):用 spec.note 当副标题,或者空
            sub = spec.note or ""
            set_placeholder_text(find_placeholder(slide, mp["subtitle_idx"]), sub, size_pt=14, bold=False)

        # 结构化分段:3 个 bullet → 3 个分段卡
        if "seg_title_idxs" in mp and spec.bullets:
            # 1) 3 个分段小标题:用前 3 个 bullet 的"前缀词"(粗略切)
            #    简单做法:第 i 个 bullet 切前 6 字符当小标题,后面是 body
            for i, idx in enumerate(mp["seg_title_idxs"]):
                if i < len(spec.bullets):
                    bullet = spec.bullets[i]
                    # 简单切:取第一个"、/:/."分隔符前的部分当小标题
                    seg_title = bullet
                    for sep in ["、", ":", "：", ".", "。", ",", "，"]:
                        if sep in bullet:
                            seg_title = bullet.split(sep, 1)[0]
                            break
                    if len(seg_title) > 8:  # 太长就截前 6 字符
                        seg_title = bullet[:6]
                    set_placeholder_text(find_placeholder(slide, idx), seg_title, size_pt=16, bold=True)
                else:
                    set_placeholder_text(find_placeholder(slide, idx), "", size_pt=16, bold=False)
            for i, idx in enumerate(mp["seg_body_idxs"]):
                if i < len(spec.bullets):
                    set_placeholder_text(find_placeholder(slide, idx), spec.bullets[i], size_pt=14, bold=False)
                else:
                    set_placeholder_text(find_placeholder(slide, idx), "", size_pt=14, bold=False)
        elif "seg_body_idxs" not in mp and spec.bullets:
            # 非结构化版式:找一个最大 body 框塞入(简单兜底)
            body_ph = max(
                (ph for ph in slide.placeholders if ph.has_text_frame and ph.placeholder_format.idx not in (mp["main_title_idx"], mp.get("subtitle_idx"))),
                key=lambda ph: ph.width * ph.height if ph.width and ph.height else 0,
                default=None,
            )
            if body_ph is not None:
                tf = body_ph.text_frame
                tf.text = ""
                for i, bullet in enumerate(spec.bullets):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.line_spacing = 1.3
                    run = p.add_run()
                    run.text = "• " + bullet
                    run.font.name = "Microsoft YaHei"
                    run.font.size = Pt(16)
                    if run.font.color and run.font.color.type is None:
                        from pptx.dml.color import RGBColor
                        run.font.color.rgb = RGBColor(0x1F, 0x1F, 0x1F)

    elif spec.kind == "closing":
        # 封底（自定义版式 = layout111，模板示例 slide8 风格）
        # 模板结构（2026-06-18 实测）：
        #   idx=11  文本占位符  "北京兴华集团" 模板默认 pos=(3.94, 3.01)in  size=5.44×0.89in
        #         默认 40pt Microsoft YaHei bold, RGB 1F1F1F
        #   idx=10  页码占位符  (SLIDE_NUMBER)  模板自动填
        # 模板自带装饰（layout 层，不需代码触碰）：
        #   - 北京建筑背景图（图片 2,13.33×7.50in）
        #   - 顶部右上角红色波浪装饰
        #   - "诚信铸就品牌 专业创造价值" 14pt 文本框  pos=(3.44, 3.68)in  size=6.02×0.45in
        # 模板示例 slide8 额外添加（slide 层，**代码需要补**）：
        #   - 左上角"北京兴华 XH GLOBAL" logo  pos=(0.94, 0.82)in  size=(1.98, 0.51)in
        #   - 左下角联系方式 textbox  pos=(0.94, 5.99)in  size=8.0×0.9in  含电话/传真/地址
        #
        # **关键规范（2026-06-18 v3 修订）**：
        #   1) 标题占位符 idx=11 模板默认位置 (3.94, 3.01) 与 slogan 框 (3.44, 3.68)
        #      存在 0.22in 视觉重叠，必须**重写 spPr/xfrm 把标题上移到 y=2.50in**，
        #      使标题底边位于 3.39in，与 slogan 顶部 3.68in 保持 0.29in 间距。
        #   2) 联系方式 3 行文本统一 **16pt**（用户规范，2026-06-18 v3 修订）
        from pptx.util import Inches, Emu
        from pptx.dml.color import RGBColor
        # ---- 写主标题，并重写占位符 spPr/xfrm 上移到 y=2.50in（保持 0.29in 间距）----
        title_ph = find_placeholder(slide, mp["main_title_idx"])
        set_placeholder_text(title_ph, spec.title, size_pt=40, bold=True)
        # 重写 spPr/xfrm：把 idx=11 占位符移到 (3.94, 2.50)in, size 5.44×0.89in
        from lxml import etree as _et
        from pptx.oxml.ns import qn
        # 找到占位符的 <p:sp> 元素
        sp_elem = title_ph._element
        spPr = sp_elem.find(qn('p:spPr'))
        if spPr is None:
            spPr = _et.SubElement(sp_elem, qn('p:spPr'))
        # 删除现有 xfrm
        for old in spPr.findall(qn('a:xfrm')):
            spPr.remove(old)
        # 新 xfrm: 位置 (3.94, 2.50)in, 尺寸 5.44×0.89in
        new_xfrm = _et.SubElement(spPr, qn('a:xfrm'))
        new_off = _et.SubElement(new_xfrm, qn('a:off'))
        new_off.set('x', str(int(3.94 * 914400)))
        new_off.set('y', str(int(2.50 * 914400)))
        new_ext = _et.SubElement(new_xfrm, qn('a:ext'))
        new_ext.set('cx', str(int(5.44 * 914400)))
        new_ext.set('cy', str(int(0.89 * 914400)))
        # ---- 补：左下角联系方式文本框（slide 层独立添加，匹配模板示例 slide8）----
        tb = slide.shapes.add_textbox(Inches(0.94), Inches(5.99), Inches(8.0), Inches(0.9))
        tb.text_frame.word_wrap = True
        lines = [
            ("电话：", "010-82250666"),
            ("传真：", "010-82250851"),
            ("地址：", "北京市西城区裕民路18号北环中心27层"),
        ]
        for i, (label, value) in enumerate(lines):
            p = tb.text_frame.paragraphs[0] if i == 0 else tb.text_frame.add_paragraph()
            r1 = p.add_run(); r1.text = label
            r1.font.name = "Microsoft YaHei"; r1.font.size = Pt(16); r1.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            r2 = p.add_run(); r2.text = value
            r2.font.name = "Microsoft YaHei"; r2.font.size = Pt(16); r2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        # ---- 补：左上角 logo（从模板 media/image2.png 抽取并插入）----
        import zipfile as _zip
        from io import BytesIO
        with _zip.ZipFile(str(Path(DEFAULT_TEMPLATE))) as z:
            logo_bytes = z.read('ppt/media/image2.png')
        slide.shapes.add_picture(BytesIO(logo_bytes), Inches(0.94), Inches(0.82),
                                  width=Inches(1.98), height=Inches(0.51))

    if spec.note and spec.kind in ("cover", "closing"):
        slide.notes_slide.notes_text_frame.text = spec.note


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("outline", help="path to outline.md")
    parser.add_argument("output", help="output .pptx path")
    parser.add_argument("template", nargs="?", default=DEFAULT_TEMPLATE, help="template .pptx path")
    parser.add_argument("--no-toc", dest="add_toc", action="store_false")
    parser.add_argument("--no-closing", dest="add_closing", action="store_false")
    args = parser.parse_args()

    outline_path = Path(args.outline)
    if not outline_path.exists():
        print(f"ERROR: outline not found: {outline_path}", file=sys.stderr)
        return 1

    template_path = Path(args.template)
    if not template_path.exists():
        print(f"ERROR: template not found: {template_path}", file=sys.stderr)
        return 1

    text = outline_path.read_text(encoding="utf-8")
    pages = parse_outline(text, add_toc=args.add_toc, add_closing=args.add_closing)

    # 关键步骤：复制模板到输出路径，再从输出路径打开。
    # 这样原始模板不被污染，输出 .pptx 自带"模板主题 + 我们写的新页"。
    # python-pptx 没有"开新文件"模式，所以走"拷贝模板 → 删除其原 slides → 加新 slides"。
    import shutil
    import tempfile
    out_path = Path(args.output)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.NamedTemporaryFile(
        suffix=".pptx", delete=False, dir=tempfile.gettempdir()
    ) as tmp:
        tmp_path = Path(tmp.name)
    shutil.copy2(str(template_path), str(tmp_path))

    prs = Presentation(str(tmp_path))

    # 删除模板中已有的全部页面（保留 master + theme + 母版 layouts）
    # 注意：只清 sldIdLst 引用还不够，必须 drop 与 slide 的 rel，否则 zip 重复。
    from pptx.oxml.ns import qn

    xml_slides = prs.slides._sldIdLst
    for sld_id in list(xml_slides):
        rId = sld_id.get(qn("r:id"))
        prs.part.drop_rel(rId)
        xml_slides.remove(sld_id)

    for spec in pages:
        render_page(prs, spec)

    prs.save(str(out_path))
    tmp_path.unlink(missing_ok=True)

    print(f"OK: wrote {out_path}")
    print(f"  pages: {len(pages)}")
    for i, spec in enumerate(pages, 1):
        print(f"  {i}. [{spec.kind:8s}] {spec.title}  ({spec.layout})")
    return 0


if __name__ == "__main__":
    sys.exit(main())
