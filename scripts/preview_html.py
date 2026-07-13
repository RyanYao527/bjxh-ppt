"""Build a self-contained HTML preview of a .pptx for user verification.

For each slide:
- shows the layout name
- lists every non-empty placeholder with idx, text, position, size, font size
- highlights obviously wrong things (text overflow, etc.)

Output: a single HTML file the user can open in any browser.

Usage:
    python preview_html.py <file.pptx>
"""

from __future__ import annotations

import html
import io
import sys
from pathlib import Path

from pptx import Presentation
from pptx.shapes.shapes import BaseShape  # noqa: F401
from pptx.util import Emu  # noqa: F401


def emu_to_in(v: int | None) -> float:
    return v / 914400 if v is not None else 0.0


def estimate_overflow(
    text: str,
    box_w_in: float,
    font_pt: int,
) -> tuple[int, int]:
    """Rough estimate: how many chars fit per line and how many lines we need.

    Returns (chars_per_line, lines_needed).

    NOTE: The CJK / English char-width ratios below are rough approximations
    (1.0 × font-size for CJK, 0.55 × for ASCII).  Actual rendering width
    depends on the font's glyph metrics and the PowerPoint text engine, so
    the overflow flag is a heuristic — always visually confirm in PowerPoint.
    """
    cjk_count: int = sum(1 for c in text if ord(c) > 0x2E80)
    en_count: int = len(text) - cjk_count
    avg_char_width_pt: float = (
        1.0 * (cjk_count / max(1, len(text)))
        + 0.55 * (en_count / max(1, len(text)))
    )
    chars_per_line: int = max(
        1, int((box_w_in * 72) / max(0.1, font_pt * avg_char_width_pt))
    )
    lines: int = max(1, -(-len(text) // chars_per_line))  # ceil division
    return chars_per_line, lines


def main() -> int:
    if len(sys.argv) < 2:
        print(__doc__)
        return 1

    src: str = sys.argv[1]
    if not Path(src).exists():
        print(f"ERROR: file not found: {src}", file=sys.stderr)
        return 1

    out_path = Path(src).with_suffix(".preview.html")

    prs = Presentation(src)
    total_slides: int = len(prs.slides)

    buf = io.StringIO()
    buf.write("<!doctype html>\n")
    buf.write('<html lang="zh-CN"><head><meta charset="utf-8">\n')
    buf.write(
        f"<title>bjxh-ppt preview: {html.escape(Path(src).name)}</title>\n"
    )
    buf.write("<style>\n")
    buf.write("""
body { font: 14px/1.5 -apple-system, "Microsoft YaHei", sans-serif; background: #f4f4f4; margin: 0; padding: 20px; }
h1 { font-size: 20px; }
h2 { font-size: 16px; margin: 0 0 8px; padding: 8px 12px; background: #1f3a5f; color: white; border-radius: 4px; }
.slide { background: white; border: 1px solid #ccc; border-radius: 6px; padding: 16px; margin-bottom: 16px; }
.slide.empty { border-color: #c33; background: #fee; }
.ph { display: grid; grid-template-columns: 50px 80px 1fr 120px 90px; gap: 8px; padding: 4px 6px; border-bottom: 1px solid #eee; font-size: 13px; align-items: start; }
.ph.idx { font-weight: bold; color: #1f3a5f; }
.ph.pos { color: #888; font-family: monospace; font-size: 12px; }
.ph.text { word-break: break-word; }
.ph.text.empty-text { color: #c33; font-style: italic; }
.ph.size { color: #666; font-size: 12px; }
.summary { background: #eef; padding: 12px; border-radius: 6px; margin-bottom: 20px; }
.summary b { color: #1f3a5f; }
.warn { color: #c33; font-weight: bold; }
.ok { color: #080; }
""")
    buf.write("</style></head><body>\n")
    buf.write("<h1>bjxh-ppt 预览报告</h1>\n")
    buf.write(f"<p>源文件: <code>{html.escape(src)}</code></p>\n")
    buf.write(
        f"<p>画布: {prs.slide_width / 914400:.3f}\" "
        f"× {prs.slide_height / 914400:.3f}\" "
        f"· 总页数: {total_slides}</p>\n"
    )

    empty_count: int = 0
    for i, slide in enumerate(prs.slides, 1):
        layout_name: str = slide.slide_layout.name
        buf.write('<div class="slide">\n')
        buf.write(
            f'<h2>slide {i} · layout: {html.escape(layout_name)}</h2>\n'
        )
        any_text: bool = False
        for ph in slide.placeholders:
            if not ph.has_text_frame:
                continue
            text: str = ph.text_frame.text
            stripped: str = text.strip()
            if stripped:
                any_text = True
            idx: int = ph.placeholder_format.idx
            pos: str = (
                f"({emu_to_in(ph.left):.2f}, {emu_to_in(ph.top):.2f})"
            )
            size: str = (
                f"{emu_to_in(ph.width):.2f}×{emu_to_in(ph.height):.2f}"
            )
            # font size: take from first run
            font_pt: int = 16
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size is not None:
                        font_pt = int(run.font.size.pt)
                        break
                if font_pt != 16:
                    break

            # overflow check
            warn: str = ""
            if stripped:
                cpl, lines = estimate_overflow(
                    stripped, emu_to_in(ph.width), font_pt
                )
                line_capacity: int = max(
                    1,
                    int(
                        (emu_to_in(ph.height) * 72)
                        / max(0.1, font_pt * 1.4)
                    ),
                )
                if lines > line_capacity:
                    warn = (
                        f' <span class="warn">[溢出! 需要{lines}行,'
                        f'但框只容{line_capacity}行]</span>'
                    )

            css: str = "ph text empty-text" if not stripped else "ph text"
            buf.write(
                f'<div class="ph"><div class="ph idx">idx={idx}</div>'
                f'<div class="ph pos">{pos}</div>'
                f'<div class="{css}">{html.escape(text) if text else "(空)"}</div>'
                f'<div class="ph size">{size}″</div>'
                f'<div class="ph size">{font_pt:.0f}pt{warn}</div></div>\n'
            )
        if not any_text:
            empty_count += 1
            buf.write('<p class="warn">⚠️ 该页所有占位符均为空</p>\n')
        buf.write("</div>\n")

    summary_ok: str = '<span class="ok">全部有内容 ✅</span>'
    summary_bad: str = f'<span class="warn">{empty_count} 个空页需要修</span>'
    summary_tail: str = summary_ok if empty_count == 0 else summary_bad
    buf.write(
        f'<div class="summary"><b>汇总</b>: 共 {total_slides} 页 · '
        f'空页 {empty_count} 个 · {summary_tail}</div>\n'
    )
    buf.write("</body></html>\n")

    out_path.write_text(buf.getvalue(), encoding="utf-8")
    print(f"wrote {out_path}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
