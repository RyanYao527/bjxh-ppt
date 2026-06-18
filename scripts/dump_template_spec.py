"""
dump_template_spec.py — read the master PPT template, dump its real spec to JSON.

This is a one-time diagnostic used to verify the SKILL.md spec numbers against
the actual master file. NOT shipped to end users; called only during skill
bootstrapping.

Usage:
    python dump_template_spec.py <path-to-template.pptx> <output.json>
"""
import json
import sys
from collections import Counter
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu


def emu_to_inches(emu: int) -> float:
    return round(emu / 914400, 3)


def main(template_path: str, output_path: str) -> None:
    prs = Presentation(template_path)

    # 1) Canvas
    canvas = {
        "width_emu": prs.slide_width,
        "height_emu": prs.slide_height,
        "width_in": emu_to_inches(prs.slide_width),
        "height_in": emu_to_inches(prs.slide_height),
    }

    # 2) Slide masters + their slide layouts
    masters = []
    for m_idx, master in enumerate(prs.slide_masters):
        layouts = []
        for l_idx, layout in enumerate(master.slide_layouts):
            layouts.append(
                {
                    "index": l_idx,
                    "name": layout.name,
                    "placeholder_count": len(layout.placeholders),
                }
            )
        masters.append(
            {
                "index": m_idx,
                "layout_count": len(master.slide_layouts),
                "layouts": layouts,
            }
        )

    # 3) Aggregate layout name frequency across all slides (so we know which
    #    layouts the template actually uses in its demo pages)
    used_layouts: Counter = Counter()
    slide_summaries = []
    for s_idx, slide in enumerate(prs.slides):
        used_layouts[slide.slide_layout.name] += 1
        if s_idx < 30:  # sample first 30 slides for the dump
            slide_summaries.append(
                {
                    "index": s_idx,
                    "layout": slide.slide_layout.name,
                    "shape_count": len(slide.shapes),
                    "has_chart": any(
                        shape.has_chart for shape in slide.shapes
                    ),
                    "has_table": any(
                        shape.has_table for shape in slide.shapes
                    ),
                }
            )

    # 4) Scan first 60 slides for font/size/color sample (real usage, not layout)
    font_counter: Counter = Counter()
    size_counter: Counter = Counter()
    for slide in list(prs.slides)[:60]:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    if run.font.name:
                        font_counter[run.font.name] += 1
                    if run.font.size:
                        size_counter[run.font.size.pt] += 1

    spec = {
        "source": str(Path(template_path).resolve()),
        "slide_count": len(prs.slides),
        "canvas": canvas,
        "masters": masters,
        "used_layouts_top": used_layouts.most_common(20),
        "sample_slides": slide_summaries,
        "font_sample": font_counter.most_common(10),
        "size_sample_pt": size_counter.most_common(15),
    }

    Path(output_path).write_text(
        json.dumps(spec, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    print(f"Wrote {output_path}")
    print(f"Slides: {spec['slide_count']}, Canvas: {canvas['width_in']}\" x {canvas['height_in']}\"")
    print(f"Top layouts: {spec['used_layouts_top'][:5]}")
    print(f"Top fonts: {spec['font_sample'][:5]}")
    print(f"Top sizes: {spec['size_sample_pt'][:5]}")


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print(__doc__)
        sys.exit(1)
    main(sys.argv[1], sys.argv[2])
