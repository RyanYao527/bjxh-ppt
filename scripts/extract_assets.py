"""
extract_assets.py — Extract reusable assets from a Beijing Xinghua PPT material deck.

Usage:
    python scripts/extract_assets.py <source.pptx> <output_dir>

Outputs:
    output_dir/
        renders/          one PNG per slide (preview)
        images/           extracted original pictures from PPT
        tables/           CSV + JSON for each table
        text-snippets.md  text content per slide
        index.json        machine-readable index
        README.md         human-readable index

Requires:
    - LibreOffice (soffice) on PATH or at tools/libreoffice/
    - PyMuPDF (fitz) for PDF -> PNG
    - python-pptx for shape extraction
"""

from __future__ import annotations

import argparse
import csv
import json
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

from pptx import Presentation


def find_soffice() -> str | None:
    """Find soffice executable."""
    soffice: str | None = shutil.which("soffice")
    if soffice:
        return soffice
    # fallback to bundled LibreOffice under tools/
    repo_root = Path(__file__).parent.parent
    bundled = (
        repo_root
        / "tools"
        / "libreoffice"
        / "SourceDir"
        / "LibreOffice"
        / "program"
        / "soffice.exe"
    )
    if bundled.exists():
        return str(bundled)
    return None


def pptx_to_pdf(pptx_path: Path, out_dir: Path, soffice: str) -> Path:
    """Convert pptx to pdf using LibreOffice headless."""
    cmd: list[str] = [
        soffice,
        "--headless",
        "--convert-to", "pdf",
        "--outdir", str(out_dir),
        str(pptx_path),
    ]
    subprocess.run(
        cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.STDOUT
    )
    pdf_path: Path = out_dir / (pptx_path.stem + ".pdf")
    if not pdf_path.exists():
        raise FileNotFoundError(
            f"LibreOffice did not produce expected PDF: {pdf_path}"
        )
    return pdf_path


def pdf_to_pngs(
    pdf_path: Path, out_dir: Path, dpi: int = 150
) -> list[Path]:
    """Render each PDF page to a PNG using PyMuPDF."""
    import fitz

    out_dir.mkdir(parents=True, exist_ok=True)
    doc = fitz.open(pdf_path)
    paths: list[Path] = []
    for i, page in enumerate(doc, 1):
        pix = page.get_pixmap(dpi=dpi)
        out_path: Path = out_dir / f"slide_{i:03d}.png"
        pix.save(out_path)
        paths.append(out_path)
    doc.close()
    return paths


def extract_images(prs: Presentation, out_dir: Path) -> list[dict[str, object]]:
    """Extract embedded pictures from PPT."""
    out_dir.mkdir(parents=True, exist_ok=True)
    extracted: list[dict[str, object]] = []
    img_counter: int = 0
    for s_idx, slide in enumerate(prs.slides, 1):
        for shp in slide.shapes:
            if shp.shape_type == 13:  # PICTURE
                img_counter += 1
                try:
                    img = shp.image
                    ext: str = img.ext or "png"
                    blob: bytes = img.blob
                    size: tuple[int, int] | None = img.size
                except (AttributeError, KeyError, OSError) as e:
                    print(
                        f"  Warning: cannot read image on slide {s_idx}: {e}",
                        file=sys.stderr,
                    )
                    continue
                out_path: Path = (
                    out_dir / f"slide_{s_idx:03d}_img_{img_counter:02d}.{ext}"
                )
                out_path.write_bytes(blob)
                extracted.append(
                    {
                        "slide": s_idx,
                        "file": out_path.name,
                        "width_px": size[0] if size else None,
                        "height_px": size[1] if size else None,
                    }
                )
    return extracted


def table_to_data(table: Table) -> list[list[str]]:
    """Convert python-pptx Table to 2D list of strings."""
    data: list[list[str]] = []
    for row in table.rows:
        row_data: list[str] = []
        for cell in row.cells:
            row_data.append(cell.text.strip().replace("\n", " "))
        data.append(row_data)
    return data


def extract_tables(prs: Presentation, out_dir: Path) -> list[dict[str, object]]:
    """Extract tables as CSV + JSON."""
    out_dir.mkdir(parents=True, exist_ok=True)
    extracted: list[dict[str, object]] = []
    for s_idx, slide in enumerate(prs.slides, 1):
        t_idx: int = 0
        for shp in slide.shapes:
            if shp.has_table:
                t_idx += 1
                data: list[list[str]] = table_to_data(shp.table)
                base_name: str = f"slide_{s_idx:03d}_table_{t_idx:02d}"

                # CSV
                csv_path: Path = out_dir / f"{base_name}.csv"
                with csv_path.open("w", newline="", encoding="utf-8-sig") as f:
                    writer = csv.writer(f)
                    writer.writerows(data)

                # JSON
                json_path: Path = out_dir / f"{base_name}.json"
                json_path.write_text(
                    json.dumps(
                        {"slide": s_idx, "table_index": t_idx, "data": data},
                        ensure_ascii=False,
                        indent=2,
                    ),
                    encoding="utf-8",
                )

                extracted.append(
                    {
                        "slide": s_idx,
                        "table_index": t_idx,
                        "csv": csv_path.name,
                        "json": json_path.name,
                        "rows": len(data),
                        "cols": len(data[0]) if data else 0,
                    }
                )
    return extracted


def extract_text_snippets(
    prs: Presentation, out_dir: Path
) -> list[dict[str, object]]:
    """Extract text content per slide to a Markdown file."""
    out_dir.mkdir(parents=True, exist_ok=True)
    md_path: Path = out_dir / "text-snippets.md"
    snippets: list[dict[str, object]] = []
    with md_path.open("w", encoding="utf-8") as f:
        f.write("# 北京兴华 PPT 素材文字片段\n\n")
        for s_idx, slide in enumerate(prs.slides, 1):
            layout: str = slide.slide_layout.name
            texts: list[str] = []
            for shp in slide.shapes:
                try:
                    t = shp.text.strip()
                    if t:
                        texts.append(t)
                except (AttributeError, TypeError):
                    pass
            # Deduplicate while preserving order
            seen: set[str] = set()
            unique_texts: list[str] = []
            for t in texts:
                if t not in seen:
                    seen.add(t)
                    unique_texts.append(t)

            snippets.append(
                {"slide": s_idx, "layout": layout, "texts": unique_texts}
            )
            f.write(f"## Slide {s_idx} ({layout})\n\n")
            if unique_texts:
                for t in unique_texts:
                    safe: str = t.replace("\n", "  \n")
                    f.write(f"- {safe}\n")
            else:
                f.write("_（无文本内容）_\n")
            f.write("\n")
    return snippets


def build_index(
    output_dir: Path,
    render_paths: list[Path],
    images: list[dict[str, object]],
    tables: list[dict[str, object]],
    snippets: list[dict[str, object]],
    source_name: str,
) -> None:
    """Create assets/index.json and assets/README.md."""
    slides: list[dict[str, object]] = []
    for s_idx, render_path in enumerate(render_paths, 1):
        slide_images: list[dict[str, object]] = [
            img for img in images if img["slide"] == s_idx
        ]
        slide_tables: list[dict[str, object]] = [
            t for t in tables if t["slide"] == s_idx
        ]
        slide_texts: list[str] = next(
            (
                s["texts"]
                for s in snippets
                if isinstance(s["texts"], list) and s["slide"] == s_idx
            ),
            [],
        )  # type: ignore[assignment]
        layout: str = next(
            (
                s["layout"]
                for s in snippets
                if isinstance(s["layout"], str) and s["slide"] == s_idx
            ),
            "",
        )  # type: ignore[assignment]

        tags: list[str] = []
        if slide_images:
            tags.append("image")
        if slide_tables:
            tags.append("table")
        if slide_texts:
            tags.append("text")

        slides.append(
            {
                "slide": s_idx,
                "layout": layout,
                "tags": tags,
                "render": f"renders/{render_path.name}",
                "images": [
                    f"images/{img['file']}" for img in slide_images
                ],
                "tables": [
                    {"csv": f"tables/{t['csv']}", "json": f"tables/{t['json']}"}
                    for t in slide_tables
                ],
                "texts": slide_texts[:10],  # preview only
            }
        )

    index: dict[str, object] = {
        "source": source_name,
        "total_slides": len(render_paths),
        "images_count": len(images),
        "tables_count": len(tables),
        "slides": slides,
    }

    index_path: Path = output_dir / "index.json"
    index_path.write_text(
        json.dumps(index, ensure_ascii=False, indent=2), encoding="utf-8"
    )

    # README.md
    readme: Path = output_dir / "README.md"
    with readme.open("w", encoding="utf-8") as f:
        f.write("# 北京兴华 PPT 素材库\n\n")
        f.write(f"源文件：`{source_name}`\n\n")
        f.write("## 素材概览\n\n")
        f.write(f"- 总页数：{len(render_paths)}\n")
        f.write(f"- 提取图片：{len(images)} 张\n")
        f.write(f"- 提取表格：{len(tables)} 个\n")
        f.write("\n")
        f.write("## 目录说明\n\n")
        f.write(
            "- `renders/`：每页 PPT 的 PNG 预览（150 DPI），用于快速浏览素材。\n"
        )
        f.write("- `images/`：从 PPT 中直接提取的原始图片文件。\n")
        f.write("- `tables/`：表格数据，每个表格一个 CSV 和一个 JSON。\n")
        f.write("- `text-snippets.md`：每页的文字内容汇总，便于复制使用。\n")
        f.write("- `index.json`：机器可读的素材索引。\n")
        f.write("\n")
        f.write("## 使用方式\n\n")
        f.write("1. 在 `renders/` 中浏览找到需要的素材页。\n")
        f.write("2. 若需要原始图片，到 `images/` 查找对应 slide 的文件。\n")
        f.write("3. 若需要表格数据，到 `tables/` 查找 CSV/JSON。\n")
        f.write("4. 若需要文字内容，直接参考 `text-snippets.md`。\n")
        f.write("\n")
        f.write("## 素材清单\n\n")
        f.write("| Slide | 版式 | 标签 | 说明 |\n")
        f.write("|------|------|------|------|\n")
        for s in slides:
            tag_str: str = ", ".join(
                s["tags"] if isinstance(s["tags"], list) else []
            )
            hint: str = " ".join(
                s["texts"][:3]
                if isinstance(s.get("texts"), list)
                else []
            )[:40]
            f.write(
                f"| {s['slide']} | {s['layout']} | {tag_str} | {hint} |\n"
            )


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("source", help="Source .pptx file")
    parser.add_argument("output", help="Output directory for assets")
    parser.add_argument(
        "--dpi", type=int, default=150,
        help="DPI for PNG rendering (default 150)",
    )
    args = parser.parse_args()

    source_path = Path(args.source)
    output_dir = Path(args.output)
    if not source_path.exists():
        print(f"ERROR: source not found: {source_path}", file=sys.stderr)
        return 1

    soffice: str | None = find_soffice()
    if not soffice:
        print(
            "ERROR: soffice not found. "
            "Please install LibreOffice or place it under tools/libreoffice/",
            file=sys.stderr,
        )
        return 1
    print(f"Using soffice: {soffice}")

    output_dir.mkdir(parents=True, exist_ok=True)

    with tempfile.TemporaryDirectory() as tmp:
        tmp_path = Path(tmp)
        print("Converting PPTX to PDF...")
        pdf_path: Path = pptx_to_pdf(source_path, tmp_path, soffice)
        print(f"PDF: {pdf_path}")

        print(f"Rendering PDF to PNG ({args.dpi} DPI)...")
        render_paths: list[Path] = pdf_to_pngs(
            pdf_path, output_dir / "renders", dpi=args.dpi
        )
        print(f"Rendered {len(render_paths)} PNGs")

    print("Loading PPTX with python-pptx...")
    prs = Presentation(str(source_path))

    print("Extracting images...")
    images: list[dict[str, object]] = extract_images(
        prs, output_dir / "images"
    )
    print(f"Extracted {len(images)} images")

    print("Extracting tables...")
    tables: list[dict[str, object]] = extract_tables(
        prs, output_dir / "tables"
    )
    print(f"Extracted {len(tables)} tables")

    print("Extracting text snippets...")
    snippets: list[dict[str, object]] = extract_text_snippets(
        prs, output_dir
    )
    print(f"Extracted text from {len(snippets)} slides")

    print("Building index...")
    build_index(
        output_dir, render_paths, images, tables, snippets, source_path.name
    )

    print(f"\nDone. Assets written to: {output_dir}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
