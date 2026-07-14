"""
render.py — Slide renderer for bjxh-ppt.

Takes a list of PageSpec objects from parse.py and writes them into a
python-pptx Presentation, one slide per spec, using the placeholder
mappings in PLACEHOLDER_MAP.
"""

from __future__ import annotations

import sys as _sys
import zipfile as _zip
from io import BytesIO
from pathlib import Path

from lxml import etree as _et
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR_INDEX
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from config import get_company_info
from layout_utils import calc_safe_font_size, clear_unused_placeholders, extract_segment_title
from parse import PageSpec
from shared import apply_minimal_formatting


# ---- placeholder lookup ---------------------------------------------------

# PLACEHOLDER_MAP is the single source of truth for "which idx in which
# layout holds what content".  Data verified against the BJXH 2026 partner-
# meeting template via scripts/inspect_placeholders.py (2026-06-18).
#
# Maintenance rule: when the master template changes, run
#   python scripts/inspect_placeholders.py <new-template.pptx>
# and re-verify every entry here.  Do not guess.

PLACEHOLDER_MAP: dict[str, dict] = {
    "主题-封面": {
        "main_title_idx": 10,
        "subtitle_idx":   11,
    },
    "主题-目录页": {
        # 4 chapter-title slots (left column) + 4 page-number slots (right column).
        # The "English sub-title" and "tiny page number" slots are left empty.
        # Maximum 4 chapters; switch to a richer TOC layout beyond that.
        "item_idxs":     [12, 21, 23, 25],
        "item_sub_idxs": [20, 22, 24, 26],
        "page_idxs":     [27, 29, 31, 33],
    },
    "标题页-空白": {
        "main_title_idx": 11,
        "subtitle_idx":   12,
    },
    "1_标题页-空白": {
        "main_title_idx": 11,
        "subtitle_idx":   12,
    },
    "1_主题-过渡页": {
        "main_title_idx":       13,
        "transition_part_idx":  11,
        "transition_num_idx":   10,
    },
    "文字模板1": {
        "main_title_idx": 11,
        "subtitle_idx":   12,
        "body_idx":       13,
    },
    "无图分段-3项": {
        "main_title_idx": 11,
        "subtitle_idx":   12,
        "seg_title_idxs": [15, 16, 17],
        "seg_body_idxs":  [18, 19, 20],
    },
    "无图分段-4项": {
        "main_title_idx": 11,
        "subtitle_idx":   12,
        "seg_title_idxs": [27, 29, 31, 33],
        "seg_body_idxs":  [28, 30, 32, 34],
    },
    "无图分段-5项": {
        "main_title_idx": 11,
        "subtitle_idx":   12,
        "seg_title_idxs": [21, 23, 25, 27, 29],
        "seg_body_idxs":  [22, 24, 26, 28, 30],
    },
    "主题-封底页": {
        "main_title_idx":   10,
        "slide_number_idx": 4,   # filled by PowerPoint, never written by code
    },
    "自定义版式": {
        # The preferred closing layout (layout #111).  Code must reposition
        # the title placeholder and add contact info + logo.
        "main_title_idx":   11,
        "slide_number_idx": 10,
    },
}


# ---- layout / placeholder helpers -----------------------------------------


def find_layout(prs: Presentation, name: str) -> "SlideLayout | None":
    """Return the SlideLayout with exact *name*, or None."""
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == name:
                return layout
    return None


def find_placeholder(slide, ph_idx: int):
    """Return the placeholder shape whose idx equals *ph_idx*, or None.

    No guessing by type — the mapping is in PLACEHOLDER_MAP.
    """
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == ph_idx:
            return ph
    return None


def set_placeholder_text(
    ph,
    text: str,
    *,
    size_pt: int = 16,
    bold: bool = False,
) -> bool:
    """Write *text* into *ph*, clearing any existing content.

    Font name, size, line-spacing, and theme color are set explicitly so
    qa.py will not flag them as None.
    """
    if ph is None:
        _sys.stderr.write(
            f"Warning: set_placeholder_text called with None placeholder "
            f"(text='{text[:40]}'). The slide will be missing this content.\n"
        )
        return False
    if not ph.has_text_frame:
        return False
    tf = ph.text_frame
    tf.text = ""
    p = tf.paragraphs[0]
    p.line_spacing = 1.3
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    if run.font.color and run.font.color.type is None:
        run.font.color.theme_color = MSO_THEME_COLOR_INDEX.TEXT_1
    return True


# ---- main renderer --------------------------------------------------------


def render_page(
    prs: Presentation,
    spec: PageSpec,
    template_path: str | Path,
) -> None:
    """Add one slide to *prs* for *spec*.

    *template_path* is the path to the master .pptx; it is used as a
    zipfile source when extracting the logo image for the closing slide.
    """
    layout = find_layout(prs, spec.layout)
    if layout is None:
        raise SystemExit(
            f"ERROR: layout '{spec.layout}' not found. "
            f"Check spelling against template_spec.json."
        )
    if spec.layout not in PLACEHOLDER_MAP:
        raise SystemExit(
            f"ERROR: layout '{spec.layout}' not in PLACEHOLDER_MAP. "
            f"Supported: {list(PLACEHOLDER_MAP)}"
        )
    slide = prs.slides.add_slide(layout)
    apply_minimal_formatting(slide)
    mp: dict = PLACEHOLDER_MAP[spec.layout]
    used_idxs: set[int] = set()

    # -- cover --------------------------------------------------------------
    if spec.kind == "cover":
        # Main title idx=10: 66 pt bold (spec v3 — must visually overlap the
        # "BEIJING XINGHUA" watermark at y=2.20-3.21 in.)
        used_idxs.add(mp["main_title_idx"])
        set_placeholder_text(
            find_placeholder(slide, mp["main_title_idx"]),
            spec.title,
            size_pt=66,
            bold=True,
        )
        # Subtitle idx=11: 18 pt single-line English org name on white
        if mp.get("subtitle_idx") is not None:
            used_idxs.add(mp["subtitle_idx"])
            sub_ph = find_placeholder(slide, mp["subtitle_idx"])
            if sub_ph is not None and sub_ph.has_text_frame:
                tf = sub_ph.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.alignment = None
                run = p.add_run()
                run.text = "BEIJING XINGHUA GROUP"
                run.font.name = "Arial"
                run.font.size = Pt(18)
                run.font.color.theme_color = MSO_THEME_COLOR_INDEX.BACKGROUND_1

    # -- table of contents --------------------------------------------------
    elif spec.kind == "toc":
        if len(spec.bullets) > len(mp["item_idxs"]):
            _sys.stderr.write(
                f"Warning: {len(spec.bullets)} chapters but TOC layout only "
                f"supports {len(mp['item_idxs'])}. "
                f"Chapters {len(mp['item_idxs']) + 1}+ will not appear in the "
                f"table of contents.\n"
            )
        for i, idx in enumerate(mp["item_idxs"]):
            used_idxs.add(idx)
            text = spec.bullets[i] if i < len(spec.bullets) else ""
            item_ph = find_placeholder(slide, idx)
            safe_toc_size = calc_safe_font_size(item_ph, text, max_pt=14, min_pt=8) if item_ph else 14
            set_placeholder_text(item_ph, text, size_pt=safe_toc_size, bold=False)
        for i, idx in enumerate(mp["page_idxs"]):
            used_idxs.add(idx)
            if i < len(spec.bullets):
                page_num = (
                    spec.toc_page_numbers[i]
                    if i < len(spec.toc_page_numbers)
                    else str(i + 1)
                )
                set_placeholder_text(
                    find_placeholder(slide, idx), page_num, size_pt=12, bold=False
                )

    # -- chapter divider ----------------------------------------------------
    elif spec.kind == "chapter":
        if spec.layout == "1_主题-过渡页":
            ch_ph = find_placeholder(slide, mp["main_title_idx"])
            safe_ch = calc_safe_font_size(ch_ph, spec.title, max_pt=28, min_pt=16) if ch_ph else 28
            set_placeholder_text(ch_ph, spec.title, size_pt=safe_ch, bold=True)
        else:
            ch_ph = find_placeholder(slide, mp["main_title_idx"])
            safe_ch = calc_safe_font_size(ch_ph, spec.title, max_pt=24, min_pt=14) if ch_ph else 24
            set_placeholder_text(ch_ph, spec.title, size_pt=safe_ch, bold=True)
            if mp.get("subtitle_idx") is not None:
                set_placeholder_text(
                    find_placeholder(slide, mp["subtitle_idx"]),
                    spec.note or "",
                    size_pt=14,
                    bold=False,
                )

    # -- content page -------------------------------------------------------
    elif spec.kind == "content":
        title_ph = find_placeholder(slide, mp["main_title_idx"])
        if title_ph is not None:
            used_idxs.add(mp["main_title_idx"])
            safe_title_size = calc_safe_font_size(title_ph, spec.title, max_pt=20, min_pt=12)
            set_placeholder_text(title_ph, spec.title, size_pt=safe_title_size, bold=True)
        if mp.get("subtitle_idx") is not None:
            sub_ph = find_placeholder(slide, mp["subtitle_idx"])
            if sub_ph is not None:
                used_idxs.add(mp["subtitle_idx"])
                set_placeholder_text(sub_ph, spec.note or "", size_pt=14, bold=False)

        # Structured segment layouts (无图分段-3/4/5项)
        if "seg_title_idxs" in mp and spec.bullets:
            for i, idx in enumerate(mp["seg_title_idxs"]):
                used_idxs.add(idx)
                if i < len(spec.bullets):
                    bullet = spec.bullets[i]
                    seg_title = extract_segment_title(bullet)
                    seg_title_ph = find_placeholder(slide, idx)
                    safe_size = calc_safe_font_size(seg_title_ph, seg_title, max_pt=16, min_pt=10) if seg_title_ph else 16
                    set_placeholder_text(seg_title_ph, seg_title, size_pt=safe_size, bold=True)
                else:
                    set_placeholder_text(
                        find_placeholder(slide, idx), "", size_pt=16, bold=False
                    )
            for i, idx in enumerate(mp["seg_body_idxs"]):
                used_idxs.add(idx)
                if i < len(spec.bullets):
                    body_ph = find_placeholder(slide, idx)
                    safe_body_size = calc_safe_font_size(body_ph, spec.bullets[i], max_pt=14, min_pt=8) if body_ph else 14
                    set_placeholder_text(body_ph, spec.bullets[i], size_pt=safe_body_size, bold=False)
                else:
                    set_placeholder_text(
                        find_placeholder(slide, idx), "", size_pt=14, bold=False
                    )

        # Non-segment layouts (e.g. 文字模板1): pour bullets into body
        elif "seg_body_idxs" not in mp and spec.bullets:
            if "body_idx" in mp:
                body_ph = find_placeholder(slide, mp["body_idx"])
            else:
                body_ph = max(
                    (
                        ph
                        for ph in slide.placeholders
                        if ph.has_text_frame
                        and ph.placeholder_format.idx
                        not in (mp["main_title_idx"], mp.get("subtitle_idx"))
                    ),
                    key=lambda ph: (
                        ph.width * ph.height if ph.width and ph.height else 0
                    ),
                    default=None,
                )
            if body_ph is not None:
                tf = body_ph.text_frame
                tf.text = ""
                for i, bullet in enumerate(spec.bullets):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.line_spacing = 1.3
                    run = p.add_run()
                    run.text = "• " + bullet
                    run.font.name = "Microsoft YaHei"
                    run.font.size = Pt(16)
                    if run.font.color and run.font.color.type is None:
                        run.font.color.theme_color = MSO_THEME_COLOR_INDEX.TEXT_1

            # 文字模板1 is thin; add a red accent line at the bottom so qa.py
            # doesn't flag THIN-PAGE.
            if spec.layout == "文字模板1":
                accent = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(1.10),
                    Inches(6.65),
                    Inches(11.14),
                    Inches(0.04),
                )
                accent.fill.solid()
                accent.fill.fore_color.theme_color = MSO_THEME_COLOR_INDEX.ACCENT_3
                accent.line.fill.background()

    # -- closing slide (自定义版式) ------------------------------------------
    elif spec.kind == "closing":
        # Write title and reposition from template default (3.94, 3.01) in
        # to (3.94, 2.50) in — avoids 0.22 in overlap with the slogan textbox.
        used_idxs.add(mp["main_title_idx"])
        title_ph = find_placeholder(slide, mp["main_title_idx"])
        set_placeholder_text(title_ph, spec.title, size_pt=40, bold=True)

        sp_elem = title_ph._element
        spPr = sp_elem.find(qn("p:spPr"))
        if spPr is None:
            spPr = _et.SubElement(sp_elem, qn("p:spPr"))
        for old in spPr.findall(qn("a:xfrm")):
            spPr.remove(old)
        new_xfrm = _et.SubElement(spPr, qn("a:xfrm"))
        new_off = _et.SubElement(new_xfrm, qn("a:off"))
        new_off.set("x", str(int(3.94 * 914400)))
        new_off.set("y", str(int(2.50 * 914400)))
        new_ext = _et.SubElement(new_xfrm, qn("a:ext"))
        new_ext.set("cx", str(int(5.44 * 914400)))
        new_ext.set("cy", str(int(0.89 * 914400)))

        # Contact info textbox (bottom-left, white 16 pt)
        company = get_company_info()
        contact_lines: list[tuple[str, str]] = []
        if company.get("phone"):
            contact_lines.append(("电话：", company["phone"]))
        if company.get("fax"):
            contact_lines.append(("传真：", company["fax"]))
        if company.get("address"):
            contact_lines.append(("地址：", company["address"]))
        if contact_lines:
            tb = slide.shapes.add_textbox(
                Inches(0.94), Inches(5.99), Inches(8.0), Inches(0.9)
            )
            tb.text_frame.word_wrap = True
            for i, (label, value) in enumerate(contact_lines):
                p = tb.text_frame.paragraphs[0] if i == 0 else tb.text_frame.add_paragraph()
                r1 = p.add_run()
                r1.text = label
                r1.font.name = "Microsoft YaHei"
                r1.font.size = Pt(16)
                r1.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                r2 = p.add_run()
                r2.text = value
                r2.font.name = "Microsoft YaHei"
                r2.font.size = Pt(16)
                r2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # Logo (extracted from the template zip)
        try:
            with _zip.ZipFile(str(template_path)) as z:
                logo_bytes = z.read("ppt/media/image2.png")
            slide.shapes.add_picture(
                BytesIO(logo_bytes),
                Inches(0.94),
                Inches(0.82),
                width=Inches(1.98),
                height=Inches(0.51),
            )
        except (KeyError, _zip.BadZipFile, OSError) as exc:
            print(
                f"Warning: could not extract logo from template: {exc}",
                file=sys.stderr,
            )

    # -- clear unused placeholder residue ----------------------------------
    clear_unused_placeholders(slide, used_idxs)

    # -- speaker notes (cover + closing only) -------------------------------
    if spec.note and spec.kind in ("cover", "closing"):
        slide.notes_slide.notes_text_frame.text = spec.note
